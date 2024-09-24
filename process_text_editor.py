from pptx import Presentation
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import RGBColor
from pdf2docx import Converter
import re
from bs4 import BeautifulSoup
import chardet
from groq import Groq
import zipfile
from lxml import etree
from io import BytesIO
import shutil
import os
import string
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
from pathlib import Path
import sys
import time
import chardet


# Script per llegir i reescriure un document (traduir o modificar el text).
# Es llegeix el text per "frases" amb format diferent, i es posa en un diccionari. Després es passen unes quantes entrades del diccionari a la IA.
# Finalment es reconstrueix el diccionari final i es torna a generar el document amb el mateix format



class Modify_Diccionarios:

    # Generar clave del diccionario para cada texto que se lee
    def generate_numeric_code(counter):
        return f"{counter}"

    # Función para no enviar a traducir textos de un color concreto, y mantenerlos como están
    def color_to_rgb(color_to_exclude):
        # Verificar si color_hex es None o está vacío
        if color_to_exclude is None or len(color_to_exclude) != 7 or not color_to_exclude.startswith('#'):
            return None
        else:
            # Si pasa la validación, convertir el color hexadecimal a RGB
            color_to_exclude = color_to_exclude.lstrip('#')
            return RGBColor(int(color_to_exclude[:2], 16), int(color_to_exclude[2:4], 16), int(color_to_exclude[4:], 16))

    # Funciones para seleccionar textos relevantes y separar el texto en bloques para enviar a traducir

    def filtrar_textos_relevantes(diccionario):
        # Filtrar textos que no estén vacíos o sólo contengan espacios en blanco
        return {code: texto for code, texto in diccionario.items() if len(texto.strip()) > 0}

    import re

    def separar_texto_bloques(diccionario, max_chars_per_block=300, min_chars_per_block=30, max_codes_per_block=5):
        """
        Separar el texto en bloques basados en un número máximo y mínimo de caracteres por bloque,
        siempre separando por el cierre de un </ph>, intentando además dividir por signos de puntuación y,
        si no es posible, por palabras en mayúscula. Limita el número de códigos <ph> por bloque a un máximo de 5.
        """

        # Generar un texto concatenado de todas las claves y valores del diccionario, encerrando el texto entre <ph> y </ph>
        full_text = "".join([f"<ph>{text}</ph>" for text in diccionario.values()])

        blocks = []
        current_pos = 0
        total_length = len(full_text)

        while current_pos < total_length:
            # Definir el límite superior e inferior para el bloque actual
            max_block_end = min(current_pos + max_chars_per_block, total_length)
            min_block_end = min(current_pos + min_chars_per_block, total_length)

            # Buscar la última posición de un </ph> antes del límite máximo
            last_placeholder_pos = max([m.end() for m in re.finditer(r'</ph>', full_text[current_pos:max_block_end])] or [0])

            # Buscar la última posición de un signo de puntuación antes del límite máximo
            last_punctuation_pos = max([m.end() for m in re.finditer(r'[.!?…]', full_text[current_pos:max_block_end])] or [0])

            # Verificar si el siguiente texto comienza con mayúscula
            next_capital_pos = re.search(r'\b[A-Z]', full_text[max_block_end:])

            # Priorizar terminar el bloque en la puntuación si existe
            if last_punctuation_pos > min_block_end:
                block_end = current_pos + last_punctuation_pos
            # Si no hay puntuación, cortar en el cierre de </ph>
            elif last_placeholder_pos > min_block_end:
                block_end = current_pos + last_placeholder_pos
            # Si no hay cierre de </ph>, cortar en la siguiente mayúscula
            elif next_capital_pos:
                block_end = max_block_end + next_capital_pos.start()
            else:
                block_end = max_block_end

            # Asegurarse de que el bloque incluya todo el contenido del último <ph>...</ph> si empieza dentro del bloque
            temp_block = full_text[current_pos:block_end]
            open_ph_pos = temp_block.rfind('<ph>')
            close_ph_pos = temp_block.rfind('</ph>')

            if open_ph_pos > close_ph_pos:
                # Si hay un <ph> abierto sin cerrar dentro del bloque, extendemos el bloque hasta encontrar el cierre
                next_close_ph = re.search(r'</ph>', full_text[block_end:])
                if next_close_ph:
                    block_end += next_close_ph.end()

            # Extraer el bloque temporalmente
            temp_block = full_text[current_pos:block_end].strip()

            # Contar el número de códigos <ph> en el bloque temporal
            num_codes = len(re.findall(r'<ph>.*?</ph>', temp_block))

            # Si el número de códigos excede el máximo permitido, ajustar el block_end
            if num_codes > max_codes_per_block:
                code_positions = [m.end() for m in re.finditer(r'</ph>', temp_block)]
                block_end = current_pos + code_positions[max_codes_per_block - 1]

            # Añadir el bloque a la lista
            block = full_text[current_pos:block_end].strip()
            blocks.append(block)

            # Actualizar la posición actual
            current_pos = block_end

        # Ajustar el penúltimo bloque si el último bloque es demasiado corto
        if len(blocks) > 1 and len(blocks[-1]) < min_chars_per_block:
            blocks[-2] += " " + blocks[-1]
            blocks.pop()

        return blocks


    # Ajustamos espacios i símbolos de cada entrada del diccionario para que sean idénticos al diccionario original
    def ajuste_post_traduccion_dict(diccionario_original, diccionario_traduccion):
        diccionario_ajustado = {}
        
        for clave, texto_original in diccionario_original.items():
            traduccion = diccionario_traduccion.get(clave)
            
            if texto_original is None or traduccion is None:
                diccionario_ajustado[clave] = traduccion
                continue

            # Extraer los espacios al inicio y final del texto original
            inicio_original_match = re.match(r'^\s*', texto_original)
            final_original_match = re.search(r'\s*$', texto_original)

            inicio_original = inicio_original_match.group(0) if inicio_original_match else ""
            final_original = final_original_match.group(0) if final_original_match else ""

            # Limpiar espacios al inicio y final de la traducción
            contenido_traduccion = traduccion.strip()

            # Concatenar los espacios del original con la traducción limpia
            texto_ajustado = inicio_original + contenido_traduccion + final_original

            diccionario_ajustado[clave] = texto_ajustado

        return diccionario_ajustado

class Modify_Bloques:
    
    # Función para juntar los bloques
    def join_blocks(bloques_traducidos):
        # Unir todos los bloques traducidos en un solo string
        traduccion_completa = "".join(bloques_traducidos)
        
        # Buscar todos los textos entre <ph> y </ph>
        translated_texts = re.findall(r"<ph>(.*?)</ph>", traduccion_completa, re.DOTALL)
        
        # Crear un diccionario con índices numéricos como claves y los textos como valores
        return {str(i): text.strip() for i, text in enumerate(translated_texts, 1)}

class Validar_Bloques:
    def verificar_placeholders(original, traducido):
        """
        Verifica que los placeholders en el texto original y traducido cumplan con las siguientes condiciones:
        1. Hay el mismo número de <ph> y </ph> en el texto original y la traducción.
        2. Un <ph> siempre debe cerrarse con un </ph>. No puede haber otro <ph> antes.
        3. Siempre debe empezar el bloque de texto por <ph> y acabar por </ph>.
        4. Los placeholders deben estar en el mismo orden en ambos textos.
        """
        def contar_y_verificar_placeholders(texto):
            apertura = texto.count('<ph>')
            cierre = texto.count('</ph>')
            if apertura != cierre:
                print(f'[ERROR] Número desigual de etiquetas de apertura y cierre: {apertura} <ph>, {cierre} </ph>')
                return False
            if not texto.startswith('<ph>') or not texto.endswith('</ph>'):
                print('[ERROR] El texto no comienza con <ph> o no termina con </ph>')
                return False
            patron = r'<ph>(?:(?!</ph>).)*?</ph>'
            bloques = re.findall(patron, texto)
            if len(bloques) != apertura:
                print('[ERROR] Algunos placeholders no están correctamente cerrados')
                return False
            return True

        # Verificar el texto original
        if not contar_y_verificar_placeholders(original):
            print('[ERROR] El texto original no cumple con los requisitos de los placeholders')
            return False

        # Verificar el texto traducido
        if not contar_y_verificar_placeholders(traducido):
            print('[ERROR] El texto traducido no cumple con los requisitos de los placeholders')
            return False

        # Verificar que los placeholders estén en el mismo orden
        placeholders_original = re.findall(r'<ph>', original)
        placeholders_traducido = re.findall(r'<ph>', traducido)
        
        if len(placeholders_original) != len(placeholders_traducido):
            print('[ERROR] El número de placeholders en la traducción no coincide con el original')
            return False

        return True

    def placeholders_por_espacios(data):
        """
        Reemplaza los placeholders <ph>...</ph> con espacios en un texto.

        :param data: Debe ser un texto (str) del cual reemplazar los placeholders.
        :return: El texto con los placeholders reemplazados por espacios.
        """
        if isinstance(data, str):
            # Reemplazamos los placeholders por espacios
            data = re.sub(r'<ph>.*?</ph>', " ", data)
            data = re.sub(r'\s+', " ", data)  # Reemplazar múltiples espacios por uno solo
            return data.strip()
        else:
            raise TypeError("La entrada para reemplazar placeholders debe ser un texto.")


    def eliminar_placeholders(data):
        """
        Elimina los placeholders <ph>...</ph> de un texto o de un diccionario.

        :param data: Puede ser un texto (str) o un diccionario (dict) del cual eliminar los placeholders.
        :return: El texto o diccionario sin los placeholders.
        """
        if isinstance(data, str):
            # Si 'data' es un texto, eliminamos los placeholders del texto
            return re.sub(r'<ph>.*?</ph>', "", data)
        
        elif isinstance(data, dict):
            # Si 'data' es un diccionario, eliminamos los placeholders para cada valor del diccionario
            return {key: re.sub(r'<ph>.*?</ph>', "", value) for key, value in data.items()}
        
        else:
            raise TypeError("La entrada para eliminar placeholders debe ser un texto o un diccionario.")

class PPTX_process:
    def correcciones_pptx(input_path, temp_output_path):
        """
        Procesa el archivo PPTX, unifica los runs con el mismo formato,
        y trata correctamente los textos en las diapositivas. Crea un archivo PPTX temporal con los cambios.
        """
        namespaces = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        def unir_runs(tree):
            """
            Unifica los 'runs' consecutivos que tienen el mismo formato en las diapositivas.
            """
            # Buscar todos los párrafos en la diapositiva
            parrafos = tree.xpath('//a:p', namespaces=namespaces)
            
            for parrafo in parrafos:
                process_paragraphs_in_container(parrafo)  # Procesar los runs dentro del párrafo

        def process_paragraphs_in_container(container):
            """
            Procesa los 'runs' dentro de un párrafo,
            unificando los que tienen el mismo formato y preservando correctamente los espacios.
            """
            runs = container.xpath('.//a:r', namespaces=namespaces)
            previous_run = None

            def es_run_vacio(run):
                """Determina si el run está vacío o tiene solo espacios o tabuladores"""
                texto = ''.join(run.xpath('.//a:t/text()', namespaces=namespaces)).strip()
                return texto == ""

            for run in runs:
                # Obtener el texto actual del run
                texto_actual = ''.join(run.xpath('.//a:t/text()', namespaces=namespaces))

                # Verificar si xml:space="preserve" está presente
                preserve_space_actual = run.xpath('.//a:t[@xml:space="preserve"]', namespaces=namespaces)

                # Eliminar propiedades irrelevantes del formato actual
                formato_run_actual = run.xpath('./a:rPr', namespaces=namespaces)
                if formato_run_actual:
                    # Elimina <a:latin>, <a:ea>, <a:cs> (equivalentes a fuentes o lenguajes) del formato actual
                    for tag in ['latin', 'ea', 'cs']:
                        for elem in formato_run_actual[0].xpath(f'.//a:{tag}', namespaces=namespaces):
                            elem.getparent().remove(elem)

                    # Convertir el formato en cadena
                    formato_run_actual_str = etree.tostring(formato_run_actual[0], method='c14n')
                else:
                    formato_run_actual_str = None

                if previous_run is not None:
                    # Obtener el texto y el formato del run anterior
                    texto_anterior = ''.join(previous_run.xpath('.//a:t/text()', namespaces=namespaces))
                    preserve_space_anterior = previous_run.xpath('.//a:t[@xml:space="preserve"]', namespaces=namespaces)

                    # Eliminar propiedades irrelevantes del formato anterior
                    formato_run_anterior = previous_run.xpath('./a:rPr', namespaces=namespaces)
                    if formato_run_anterior:
                        for tag in ['latin', 'ea', 'cs']:
                            for elem in formato_run_anterior[0].xpath(f'.//a:{tag}', namespaces=namespaces):
                                elem.getparent().remove(elem)

                        # Convertir el formato en cadena
                        formato_run_anterior_str = etree.tostring(formato_run_anterior[0], method='c14n')
                    else:
                        formato_run_anterior_str = None

                    # Unir los runs si tienen el mismo formato o si el run actual está vacío
                    if formato_run_actual_str == formato_run_anterior_str or es_run_vacio(run):
                        # Añadir un espacio si el anterior run tenía xml:space="preserve"
                        if preserve_space_anterior or preserve_space_actual:
                            texto_anterior += " " + texto_actual.strip()  # Agregar espacio si es necesario
                        else:
                            texto_anterior += texto_actual  # Unir directamente si no se necesita espacio adicional

                        # Actualizar el texto del run anterior
                        previous_run.xpath('.//a:t', namespaces=namespaces)[0].text = texto_anterior
                        run.getparent().remove(run)  # Eliminar el run actual ya que lo unimos al anterior
                    else:
                        previous_run = run
                else:
                    previous_run = run

        try:
            # Crear un archivo temporal para aplicar las correcciones
            with zipfile.ZipFile(input_path, 'r') as pptx_zip:
                temp_zip_path = temp_output_path + '_temp.zip'

                with zipfile.ZipFile(temp_zip_path, 'w') as temp_zip:
                    # Procesar cada diapositiva de ppt/slides/slideN.xml
                    for item in pptx_zip.infolist():
                        if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                            with pptx_zip.open(item.filename) as slide_xml:
                                tree = etree.parse(slide_xml)
                                unir_runs(tree)  # Unificar los runs en la diapositiva

                            # Escribir los cambios en el archivo temporal
                            with BytesIO() as buffer:
                                tree.write(buffer, xml_declaration=True, encoding='UTF-8')
                                buffer.seek(0)
                                temp_zip.writestr(item.filename, buffer.read())
                        else:
                            temp_zip.writestr(item, pptx_zip.read(item.filename))

            # Renombrar el archivo temporal como el archivo de salida final
            shutil.move(temp_zip_path, temp_output_path)
            print(f"Correcciones aplicadas y guardadas en {temp_output_path}")

        except zipfile.BadZipFile:
            print(f"El archivo {input_path} no es un archivo ZIP válido o está corrupto.")
        except FileNotFoundError:
            print(f"El archivo {input_path} no se encuentra.")
        except Exception as e:
            print(f"Ha ocurrido un error inesperado: {e}")




    # Lectura/Reemplazo de PPTX
    def leer_doc(input_path,output_path, color_to_exclude, textos_traducidos_final, action): 
                

        doc = Presentation(input_path)

        exclude_color_rgb = Modify_Diccionarios.color_to_rgb(color_to_exclude)

        textos_originales = {}
        counter = 1  # Inicializar el contador desde 1

        # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
        for slide in doc.slides:
            for shape in slide.shapes: 

                if shape.has_text_frame: # Acceder a figuras y cuadros de texto
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text: # Verificamos que existe la figura
                                if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                    code = Modify_Diccionarios.generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final and (exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb):
                                        run.text = textos_traducidos_final[code]

                elif shape.has_table: # Acceder a tablas
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text: # Verificamos que existe la figura
                                        if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                            code = Modify_Diccionarios.generate_numeric_code(counter)
                                            counter += 1
                                            if action == "leer":
                                                textos_originales[code] = run.text
                                            elif action == "reemplazar" and code in textos_traducidos_final and (exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb):
                                                    run.text = textos_traducidos_final[code]

                elif shape.has_chart:  # Acceder a gráficos
                    chart = shape.chart
                    # Título del gráfico
                    element = chart.chart_title
                    if element.has_text_frame:
                        if element.text_frame.text.strip():
                            code = Modify_Diccionarios.generate_numeric_code(counter)
                            counter += 1
                            if action == "leer":
                                textos_originales[code] = element.text_frame.text
                            elif action == "reemplazar" and code in textos_traducidos_final and (exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb):
                                element.text_frame.text = textos_traducidos_final[code]

                elif shape.has_chart:  # Acceder a gráficos
                    chart = shape.chart
                    for serie in chart:
                        element = serie.point.data_labels
                        if element.has_text_frame:
                            if element.text_frame.text.strip():
                                code = Modify_Diccionarios.generate_numeric_code(counter)
                                counter += 1
                                if action == "leer":
                                    textos_originales[code] = element.text_frame.text
                                elif action == "reemplazar" and code in textos_traducidos_final and (exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb):
                                    element.text_frame.text = textos_traducidos_final[code]

        if action == "leer":
            return textos_originales
        elif action == "reemplazar":
            return doc.save(output_path)
        
class DOCX_process:
    
    @staticmethod
    def correcciones_docx(input_path, temp_output_path):
        """
        Procesa el archivo DOCX, unifica los runs con el mismo formato,
        y trata correctamente las tablas y listas. Crea un archivo DOCX temporal con los cambios.
        """
        namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

        def unir_runs(tree):
            """
            Unifica los 'runs' consecutivos que tienen el mismo formato (tanto a nivel de run como de párrafo),
            y mantiene separados los textos en tablas y listas.
            """
            # Buscar todas las tablas
            tablas = tree.xpath('//w:tbl', namespaces=namespaces)
            
            for tabla in tablas:
                # Procesar cada celda de la tabla por separado
                celdas = tabla.xpath('.//w:tc', namespaces=namespaces)
                for celda in celdas:
                    process_paragraphs_in_container(celda)  # Procesar runs dentro de la celda

            # Buscar todas las listas (párrafos con numeración)
            listas = tree.xpath('//w:numPr', namespaces=namespaces)
            for lista in listas:
                parrafo = lista.xpath('./ancestor::w:p[1]', namespaces=namespaces)
                if parrafo:
                    process_paragraphs_in_container(parrafo[0])  # Procesar runs dentro del párrafo de la lista

            # Procesar todos los párrafos que no sean parte de tablas o listas
            parrafos = tree.xpath('//w:p[not(ancestor::w:tbl)]', namespaces=namespaces)
            for parrafo in parrafos:
                if not parrafo.xpath('.//w:numPr', namespaces=namespaces):  # Excluir listas
                    process_paragraphs_in_container(parrafo)  # Procesar los runs del párrafo

        def process_paragraphs_in_container(container):
            """
            Procesa los 'runs' dentro de un contenedor (párrafo, celda de tabla o lista),
            unificando los que tienen el mismo formato y preservando correctamente los espacios.
            """
            runs = container.xpath('.//w:r', namespaces=namespaces)
            previous_run = None

            def es_run_vacio(run):
                """Determina si el run está vacío o tiene solo espacios o tabuladores"""
                texto = ''.join(run.xpath('.//w:t/text()', namespaces=namespaces)).strip()
                return texto == ""

            for run in runs:
                # Obtener el texto actual del run
                texto_actual = ''.join(run.xpath('.//w:t/text()', namespaces=namespaces))

                # Verificar si xml:space="preserve" está presente
                preserve_space_actual = run.xpath('.//w:t[@xml:space="preserve"]', namespaces=namespaces)

                # Eliminar <w:spacing> y otras propiedades irrelevantes del formato actual
                formato_run_actual = run.xpath('./w:rPr', namespaces=namespaces)
                if formato_run_actual:
                    # Eliminar <w:spacing>, <w:lang>, <w:rFonts> del formato actual
                    for tag in ['spacing', 'lang', 'rFonts']:
                        for elem in formato_run_actual[0].xpath(f'.//w:{tag}', namespaces=namespaces):
                            elem.getparent().remove(elem)

                    # Convertir el formato en cadena
                    formato_run_actual_str = etree.tostring(formato_run_actual[0], method='c14n')
                else:
                    formato_run_actual_str = None

                if previous_run is not None:
                    # Obtener el texto y el formato del run anterior
                    texto_anterior = ''.join(previous_run.xpath('.//w:t/text()', namespaces=namespaces))
                    preserve_space_anterior = previous_run.xpath('.//w:t[@xml:space="preserve"]', namespaces=namespaces)

                    # Eliminar <w:spacing> y otras propiedades irrelevantes del formato anterior
                    formato_run_anterior = previous_run.xpath('./w:rPr', namespaces=namespaces)
                    if formato_run_anterior:
                        for tag in ['spacing', 'lang', 'rFonts']:
                            for elem in formato_run_anterior[0].xpath(f'.//w:{tag}', namespaces=namespaces):
                                elem.getparent().remove(elem)

                        # Convertir el formato en cadena
                        formato_run_anterior_str = etree.tostring(formato_run_anterior[0], method='c14n')
                    else:
                        formato_run_anterior_str = None

                    # Si el run actual está vacío, concatenarlo sin cambiar el formato de previous_run
                    if es_run_vacio(run):
                        # Si el run tiene `xml:space="preserve"`, aseguramos añadir un espacio.
                        if preserve_space_anterior or preserve_space_actual:
                            texto_anterior += " "  # Añadir un espacio si es necesario por preserve
                        else:
                            texto_anterior += texto_actual  # Concatenar sin agregar espacio

                        # Actualizar el texto del run anterior con el texto concatenado
                        previous_run.xpath('.//w:t', namespaces=namespaces)[0].text = texto_anterior
                        run.getparent().remove(run)  # Eliminar el run vacío ya que lo hemos unido
                        continue  # No actualizamos previous_run, seguimos con el siguiente run

                    # Unir los runs si tienen el mismo formato o si el run actual está vacío
                    if formato_run_actual_str == formato_run_anterior_str:
                        if preserve_space_anterior or preserve_space_actual:
                            texto_anterior += " " + texto_actual.strip()  # Añadir un espacio si es necesario
                        else:
                            texto_anterior += texto_actual  # Unir directamente si no se necesita espacio adicional

                        # Actualizar el texto del run anterior
                        previous_run.xpath('.//w:t', namespaces=namespaces)[0].text = texto_anterior
                        run.getparent().remove(run)  # Eliminar el run actual ya que lo unimos al anterior
                    else:
                        previous_run = run
                else:
                    previous_run = run

        try:
            # Crear un archivo temporal para aplicar las correcciones
            with zipfile.ZipFile(input_path, 'r') as docx_zip:
                temp_zip_path = temp_output_path + '_temp.zip'

                with zipfile.ZipFile(temp_zip_path, 'w') as temp_zip:
                    # Leer y procesar el archivo document.xml
                    with docx_zip.open('word/document.xml') as document_xml:
                        tree = etree.parse(document_xml)
                        unir_runs(tree)  # Unificar los runs respetando tablas y listas

                    # Crear el nuevo archivo .docx con los cambios aplicados
                    for item in docx_zip.infolist():
                        if item.filename == 'word/document.xml':
                            with BytesIO() as buffer:
                                tree.write(buffer, xml_declaration=True, encoding='UTF-8')
                                buffer.seek(0)
                                temp_zip.writestr(item.filename, buffer.read())
                        else:
                            temp_zip.writestr(item, docx_zip.read(item.filename))

            # Renombrar el archivo temporal como el archivo de salida final
            shutil.move(temp_zip_path, temp_output_path)
            print(f"Correcciones aplicadas y guardadas en {temp_output_path}")

        except zipfile.BadZipFile:
            print(f"El archivo {input_path} no es un archivo ZIP válido o está corrupto.")
        except FileNotFoundError:
            print(f"El archivo {input_path} no se encuentra.")
        except Exception as e:
            print(f"Ha ocurrido un error inesperado: {e}")

    @staticmethod
    def leer_doc(input_path, output_path, color_to_exclude, textos_traducidos_final, action):
        """
        Lee o reemplaza el texto de un archivo DOCX. Al leer, genera un diccionario de textos originales.
        Al reemplazar, usa el diccionario de textos traducidos para actualizar el DOCX.
        """
        textos_originales = {}
        counter = 1  # Inicializar el contador desde 1

        try:
            # Abrir y procesar el archivo DOCX
            with zipfile.ZipFile(input_path, 'r') as docx_zip:
                # Leer el archivo document.xml
                with docx_zip.open('word/document.xml') as document_xml:
                    tree = etree.parse(document_xml)
                    namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

                    if action == "leer":
                        textos = tree.xpath('//w:t', namespaces=namespaces)
                        for texto in textos:
                            texto_str = texto.text.strip() if texto.text else ''
                            code = Modify_Diccionarios.generate_numeric_code(counter)
                            counter += 1
                            textos_originales[code] = texto_str

                    elif action == "reemplazar":
                        textos = tree.xpath('//w:t', namespaces=namespaces)
                        for texto in textos:
                            code = Modify_Diccionarios.generate_numeric_code(counter)
                            counter += 1
                            if code in textos_traducidos_final:
                                texto.text = textos_traducidos_final[code]

                # Guardar los cambios si estamos reemplazando texto
                if action == "reemplazar":
                    temp_zip_path = output_path + '_temp.zip'
                    with zipfile.ZipFile(input_path, 'r') as docx_zip:
                        with zipfile.ZipFile(temp_zip_path, 'w') as temp_zip:
                            # Copiar todos los archivos originales
                            for item in docx_zip.infolist():
                                if item.filename == 'word/document.xml':
                                    with BytesIO() as buffer:
                                        tree.write(buffer, xml_declaration=True, encoding='UTF-8')
                                        buffer.seek(0)
                                        temp_zip.writestr(item.filename, buffer.read())
                                else:
                                    temp_zip.writestr(item, docx_zip.read(item.filename))

                    # Renombrar el archivo temporal como el archivo de salida final
                    shutil.move(temp_zip_path, output_path)

            if action == "leer":
                print(f"Diccionario textos_originales: {textos_originales}")
                return textos_originales
            else:
                return None

        except zipfile.BadZipFile:
            print(f"El archivo {input_path} no es un archivo ZIP válido o está corrupto.")
        except FileNotFoundError:
            print(f"El archivo {input_path} no se encuentra.")
        except Exception as e:
            print(f"Ha ocurrido un error inesperado: {e}")

class PDF_process:

    # Lectura/Reemplazo PDF
    def leer_doc(input_path, output_path, textos_originales, color_to_exclude, textos_traducidos_final, action): 
        # Conversor de PDF a WORD
        def pdf_to_word(input_path, word_path):
            # Convertir el PDF a Word
            cv = Converter(input_path)
            cv.convert(word_path, start=0, end=None)
            cv.close()
        # Convertir el PDF a Word
        word_path = input_path.replace('.pdf', '.docx')
        pdf_to_word(input_path, word_path)

        return DOCX_process.leer_doc(word_path,output_path, textos_originales, color_to_exclude, textos_traducidos_final, action)

class Excel_process:
    def leer_doc(input_path, output_path, color_to_exclude, textos_traducidos_final, action):
        wb = load_workbook(input_path, keep_vba=True)  # keep_vba=True si el archivo contiene macros
        exclude_color_rgb = Modify_Diccionarios.color_to_rgb(color_to_exclude)

        textos_originales = {}
        counter = 1

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        cell_value = cell.value.strip()
                        if cell_value:
                            code = Modify_Diccionarios.generate_numeric_code(counter)
                            counter += 1
                            cell_color_rgb = None
                            if cell.fill and isinstance(cell.fill, PatternFill):
                                cell_color_rgb = cell.fill.start_color.rgb[-6:]
                                if cell_color_rgb:
                                    cell_color_rgb = tuple(int(cell_color_rgb[i:i+2], 16) for i in (0, 2, 4))
                            if action == "leer":
                                textos_originales[code] = cell_value
                            elif action == "reemplazar" and code in textos_traducidos_final and (exclude_color_rgb is None or cell_color_rgb != exclude_color_rgb):
                                cell.value = textos_traducidos_final[code]

        if action == "leer":
            return textos_originales
        elif action == "reemplazar":
            wb.save(output_path)

    def procesar_original(dict):  # Por si tenemos que hacer ajustes específicos por tipo de documento
        return dict

    def reconstruir_original(dict_traducido, dict_original):  # Por si tenemos que hacer ajustes específicos por tipo de documento
        return dict_traducido

class TXT_process:
    
    @staticmethod
    def leer_doc(input_path, output_path, textos_traducidos_final, action):
        """
        Procesa archivos TXT, ya sea para leerlos y devolver un diccionario de textos originales,
        o para reemplazar los textos con textos traducidos.
        """
        textos_originales = {}

        # Lista de codificaciones a probar
        encodings = ['utf-8', 'latin-1', 'windows-1252']
        content = None
        
        # Detectar la codificación usando chardet
        try:
            with open(input_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                encoding_detected = result['encoding']
                print(f"Codificación detectada automáticamente: {encoding_detected}")

                # Intentamos leer con la codificación detectada
                try:
                    content = raw_data.decode(encoding_detected)
                except (UnicodeDecodeError, TypeError):
                    print(f"No se pudo decodificar con la codificación detectada: {encoding_detected}. Intentando otras codificaciones...")

        except Exception as e:
            print(f"Error al leer el archivo en modo binario: {e}")

        # Si la detección automática falla o no es precisa, probamos otras codificaciones
        if content is None:
            for encoding in encodings:
                try:
                    with open(input_path, 'r', encoding=encoding) as file:
                        content = file.read()
                        print(f"Archivo leído con éxito usando la codificación: {encoding}")
                        break  # Si tiene éxito, salimos del loop
                except UnicodeDecodeError:
                    print(f"Error al leer el archivo con la codificación: {encoding}")
                except Exception as e:
                    print(f"Ocurrió otro error al leer el archivo: {e}")

        if content is None:
            raise Exception("No se pudo leer el archivo con ninguna de las codificaciones disponibles.")
        
        # Procesar el contenido según la acción
        if action == "leer":
            lines = content.splitlines()
            counter = 1
            for line in lines:
                code = Modify_Diccionarios.generate_numeric_code(counter)
                textos_originales[code] = line.strip()
                counter += 1
            return textos_originales

        elif action == "reemplazar" and textos_traducidos_final is not None:
            # Reemplazar el contenido con los textos traducidos
            lines = content.splitlines()
            counter = 1
            with open(output_path, 'w', encoding='utf-8') as file:
                for line in lines:
                    code = Modify_Diccionarios.generate_numeric_code(counter)
                    if code in textos_traducidos_final:
                        file.write(textos_traducidos_final[code] + '\n')
                    else:
                        file.write(line + '\n')
                    counter += 1
        return textos_originales

class HTML_process:
    
    @staticmethod
    def leer_doc(input_path, output_path, textos_traducidos_final, action):
        """
        Procesa archivos HTML, ya sea para leer el contenido de texto y devolver un diccionario de textos originales,
        o para reemplazar los textos con textos traducidos.
        """
        textos_originales = {}

        # Lista de codificaciones a probar
        encodings = ['utf-8', 'latin-1', 'windows-1252']
        content = None
        
        # Detectar la codificación usando chardet
        try:
            with open(input_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                encoding_detected = result['encoding']
                print(f"Codificación detectada automáticamente: {encoding_detected}")

                # Intentamos leer con la codificación detectada
                try:
                    content = raw_data.decode(encoding_detected)
                except (UnicodeDecodeError, TypeError):
                    print(f"No se pudo decodificar con la codificación detectada: {encoding_detected}. Intentando otras codificaciones...")

        except Exception as e:
            print(f"Error al leer el archivo en modo binario: {e}")

        # Si la detección automática falla o no es precisa, probamos otras codificaciones
        if content is None:
            for encoding in encodings:
                try:
                    with open(input_path, 'r', encoding=encoding) as file:
                        content = file.read()
                        print(f"Archivo leído con éxito usando la codificación: {encoding}")
                        break  # Si tiene éxito, salimos del loop
                except UnicodeDecodeError:
                    print(f"Error al leer el archivo con la codificación: {encoding}")
                except Exception as e:
                    print(f"Ocurrió otro error al leer el archivo: {e}")

        if content is None:
            raise Exception("No se pudo leer el archivo con ninguna de las codificaciones disponibles.")
        
        try:
            # Usar BeautifulSoup para analizar el HTML
            soup = BeautifulSoup(content, 'html.parser')

            # Recorrer todos los elementos de texto en el HTML
            textos = soup.find_all(string=True)  # Busca todo el texto dentro del HTML

            if action == "leer":
                counter = 1
                for texto in textos:
                    texto_str = texto.strip()
                    if texto_str:  # Evitar procesar textos vacíos
                        code = Modify_Diccionarios.generate_numeric_code(counter)
                        counter += 1
                        textos_originales[code] = texto_str

            elif action == "reemplazar" and textos_traducidos_final is not None:
                counter = 1
                for texto in textos:
                    texto_str = texto.strip()
                    if texto_str:  # Evitar procesar textos vacíos
                        code = Modify_Diccionarios.generate_numeric_code(counter)
                        counter += 1
                        if code in textos_traducidos_final:
                            texto.replace_with(textos_traducidos_final[code])

                # Guardar el archivo HTML modificado
                with open(output_path, 'w', encoding='utf-8') as file:
                    file.write(str(soup.prettify()))  # Guardar el HTML formateado

            if action == "leer":
                print(f"Diccionario textos_originales: {textos_originales}")
                return textos_originales
            else:
                return None
        
        except FileNotFoundError:
            print(f"El archivo {input_path} no se encuentra.")
        except Exception as e:
            print(f"Ha ocurrido un error inesperado: {e}")
