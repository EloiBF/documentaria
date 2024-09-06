import os
import docx
import html2text
from PyPDF2 import PdfReader
from pptx import Presentation
import re
from langchain.memory import ConversationBufferMemory
from langchain.schema import HumanMessage, AIMessage
from langchain.prompts import PromptTemplate
from dateutil import parser
from groq import Groq



def read_document(file_path):
    """Lee el contenido del archivo especificado según su tipo y devuelve el texto y la extensión del archivo."""
    _, file_extension = os.path.splitext(file_path)
    
    print(f"Procesando archivo: {file_path}")
    
    if file_extension == '.txt':
        return read_txt(file_path), file_extension
    
    elif file_extension == '.docx':
        return read_docx(file_path), file_extension
    
    elif file_extension == '.pdf':
        return read_pdf(file_path), file_extension
    
    elif file_extension == '.pptx':
        return read_pptx(file_path), file_extension
    
    elif file_extension == '.html' or file_extension == '.htm':
        return read_html(file_path), file_extension
    
    else:
        raise ValueError(f"Tipo de archivo no soportado: {file_extension}")

def read_txt(file_path):
    """Lee el contenido de un archivo .txt."""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_docx(file_path):
    """Lee el contenido de un archivo .docx."""
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def read_pdf(file_path):
    """Lee el contenido de un archivo .pdf."""
    reader = PdfReader(file_path)
    full_text = []
    for page in reader.pages:
        full_text.append(page.extract_text())
    return '\n'.join(full_text)

def read_pptx(file_path):
    """Lee el contenido de un archivo .pptx."""
    presentation = Presentation(file_path)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)

def read_html(file_path):
    """Lee el contenido de un archivo .html o .htm."""
    with open(file_path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    return html2text.html2text(html_content)

def split_text(text, max_length=100000):
    """Divide el texto en bloques, asegurándose de que cada bloque no supere max_length caracteres."""
    blocks = []
    start = 0
    while start < len(text):
        end = start + max_length
        if end < len(text):
            # Intentamos cortar en un punto lógico, buscando el último punto o salto de línea antes del límite de caracteres.
            cut_index = text.rfind('.', start, end)
            if cut_index == -1:  # Si no hay un punto, cortamos directamente en el límite.
                cut_index = end
            else:
                cut_index += 1  # Incluir el punto en el bloque.
        else:
            cut_index = len(text)
        
        blocks.append(text[start:cut_index].strip())
        start = cut_index
    
    print(f"Dividiendo texto en {len(blocks)} bloques.")
    return blocks



def model_exctact_info(texto, prompt, respuesta_tipo, ejemplo_respuesta=None, file_type=None, model='llama-3.1-70b-versatile', api_key_file='API_KEY.txt'):
    try:
        # Inicializa el cliente de Groq
        with open(api_key_file, 'r') as fichero:
            api_key = fichero.read().strip()
        client = Groq(api_key=api_key)

        # Define el prompt base según el tipo de archivo y parámetros proporcionados
        ejemplo_texto = f"Examples of responses include: {ejemplo_respuesta}" if ejemplo_respuesta else ""
        
        if file_type in ['.pptx', '.docx', '.pdf', '.txt', '.html', None]:
            prompt_completo = f"""
            You are a data extraction assistant. Follow these rules:
            - Your task is to extract specific information from the text based on the user's instructions.
            - The type of response expected is: {respuesta_tipo}. {ejemplo_texto}
            - Provide only the extracted information as per the instructions. Do not add any extra comments or information.
            - Ensure that the response matches the expected format and content type.
            
            Prompt for extraction:
            {prompt}
            Text to analyze:
            {texto}
            """
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        # Llama a la API de Groq para extraer la información
        chat_completion = client.chat.completions.create(
                messages=[
                    {
                        "role": "user",
                        "content": prompt_completo,
                    }
                ],
                model=model
            )

        respuesta = chat_completion.choices[0].message.content.strip()

    except Exception as e:
        raise RuntimeError(f"Error during information extraction: {e}")
    
    return respuesta




def validar_respuesta(respuesta, tipo):
    """
    Valida la respuesta según el tipo de respuesta esperado.
    """
    if respuesta is None:
        print(f"Respuesta inválida: N/D (Ninguna respuesta proporcionada)")
        return False
    
    respuesta = respuesta.strip().lower()  # Normalizar respuesta
    
    if tipo == 'SI/NO':
        # Aceptar variaciones de sí/no en español e inglés
        if not re.match(r'^(sí?|si?|no?|y(es)?|n(o)?)$', respuesta, re.IGNORECASE):
            print(f"Respuesta inválida: {respuesta} (debe ser 'sí', si , 'no', 'yes', o 'no')")
            return False
        return True
    
    elif tipo == 'numérica':
        # Eliminar cualquier carácter no numérico (excepto dígitos)
        respuesta_sin_simbolos = re.sub(r'\D', '', respuesta)
        if not respuesta_sin_simbolos.isdigit():
            print(f"Respuesta inválida: {respuesta} (debe ser un número)")
            return False
        return True
    
    elif tipo == 'fecha':
        try:
            # Intentar parsear la fecha
            parser.parse(respuesta, fuzzy=True)
            return True
        except ValueError:
            print(f"Respuesta inválida: {respuesta} (no es una fecha válida)")
            return False
    
    elif tipo == 'texto libre':
        if not (isinstance(respuesta, str) and len(respuesta) > 0):
            print(f"Respuesta inválida: {respuesta} (debe ser un texto libre no vacío)")
            return False
        return True
    
    else:
        print(f"Tipo de respuesta desconocido: {tipo}")
        return False



def extract_with_retry(texto, prompt, respuesta_tipo, ejemplo_respuesta=None, file_type=None, model='llama-3.1-70b-versatile', api_key_file='API_KEY.txt', max_retries=3):
    """
    Intenta extraer información del texto con reintentos en caso de fallo.
    """
    for intento in range(max_retries):
        try:
            respuesta = model_exctact_info(
                texto=texto,
                prompt=prompt,
                respuesta_tipo=respuesta_tipo,
                ejemplo_respuesta=ejemplo_respuesta,
                file_type=file_type,
                model=model,
                api_key_file=api_key_file
            )
            return respuesta
        except Exception as e:
            print(f"Intento {intento + 1} de {max_retries} fallido: {e}")
            if intento == max_retries - 1:
                raise  # Re-lanzar la excepción si se han agotado los reintentos



def reflexionar_respuestas(respuestas, prompt, tipo_respuesta):
    """
    Usa el modelo para reflexionar sobre todas las respuestas obtenidas y determinar la más adecuada.
    Incluye el tipo de respuesta en el prompt para que el modelo considere el formato adecuado.
    """
    # Usamos un prompt para pedir al modelo que reflexione sobre las respuestas obtenidas

    print(prompt)
    print(respuestas)

    prompt_reflexion = f"""
    A continuación, se te proporcionan varias respuestas obtenidas para la pregunta '{prompt}':
    
    {respuestas}
    
    Considera que el tipo de respuesta esperado es: {tipo_respuesta}.
    Reflexiona sobre estas respuestas y genera una única respuesta final basada en el consenso o en la información más relevante.
    """
    
    # Enviar el prompt al modelo para que genere la respuesta final consolidada
    respuesta_consolidada = model_exctact_info(
        texto="",
        prompt=prompt_reflexion,
        respuesta_tipo=tipo_respuesta,
        model='llama-3.1-70b-versatile',  # Ajusta según el modelo que estás utilizando
        api_key_file='API_KEY.txt'
    )
    return respuesta_consolidada.strip()

def extract_info_from_doc(input_path, prompts, tipos_respuesta, ejemplos_respuesta=None, max_retries=3):
    """Función principal para extraer información de un documento usando múltiples prompts y tipos de respuesta."""
    try:
        # Leer el documento y obtener el tipo de archivo
        text, file_type = read_document(input_path)
        
        # Dividir el texto en bloques si es necesario
        blocks = split_text(text)
        
        # Diccionario para almacenar respuestas por cada pregunta
        respuestas_por_prompt = {prompt: [] for prompt in prompts}
        
        # Iterar por cada bloque y extraer información
        for block in blocks:
            print(f"Procesando bloque de texto...")
            
            for prompt, tipo, ejemplo in zip(prompts, tipos_respuesta, ejemplos_respuesta or []):
                # Llamada a la nueva función con reintentos
                extracted_info = extract_with_retry(
                    texto=block,
                    prompt=prompt,
                    respuesta_tipo=tipo,
                    ejemplo_respuesta=ejemplo,
                    file_type=file_type,
                    max_retries=max_retries
                )
                
                # Guardar la respuesta parcial en el diccionario
                respuestas_por_prompt[prompt].append(extracted_info)
        
        # Reflexionar sobre las respuestas y obtener una respuesta final por pregunta
        final_results = {}
        for prompt, tipo in zip(prompts, tipos_respuesta):
            respuestas = respuestas_por_prompt[prompt]
            
            # Reflexionar sobre las respuestas y generar una respuesta consolidada
            respuesta_final = reflexionar_respuestas(respuestas, prompt, tipo)
            
            # Validar la respuesta final antes de guardarla
            if validar_respuesta(respuesta_final, tipo):
                final_results[prompt] = respuesta_final
            else:
                final_results[prompt] = 'N/D'  # Si la validación falla, devolver "N/D"
        
        print(final_results)
        return final_results

    except Exception as e:
        print(f"Ocurrió un error: {e}")
        return None