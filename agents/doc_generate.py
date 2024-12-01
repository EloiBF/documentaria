import re
from docx import Document
from docx.shared import RGBColor, Pt
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor as PPTXRGBColor
from groq import Groq
from docx.enum.text import WD_BREAK
from docx.shared import Pt

class DocxProcessor:
    def __init__(self):
        self.doc = Document()
        self._create_styles()

    def _create_styles(self):
        # Crear un estilo H1 de párrafo
        h1_style = self.doc.styles.add_style('H1', 1)  # Usamos 1 para un estilo de párrafo
        h1_style.font.size = Pt(18)
        h1_style.font.bold = True

        # Crear un estilo H2 de párrafo
        h2_style = self.doc.styles.add_style('H2', 1)  # Usamos 1 para un estilo de párrafo
        h2_style.font.size = Pt(16)
        h2_style.font.bold = True

    def create_document(self, content, output_path='output.docx'):
        self._process_content(content)
        self.doc.save(output_path)
        return output_path

    def _process_content(self, content):
        lines = content.splitlines()
        current_table = None
        for line in lines:
            if "[TABLE]" in line:
                current_table = []
            elif "[/TABLE]" in line:
                if current_table:
                    self._create_table(current_table)
                current_table = None
            elif current_table is not None:
                current_table.append(line)
            else:
                self._process_line(line)

    def _process_line(self, line):
        # Si encontramos el marcador de salto de página, agregamos un salto de página
        if "[PAGEBREAK]" in line:
            self.doc.add_page_break()
            line = line.replace("[PAGEBREAK]", "")  # Limpiar el texto del salto de página
        
        # Añadir el párrafo y aplicar formato
        p = self.doc.add_paragraph()
        self._process_formatting(p, line)

    def _get_color(self, color_name):
        # Mapa de colores por nombre, devuelve instancias de RGBColor directamente
        color_mapping = {
            'red': RGBColor(0xFF, 0x00, 0x00),
            'green': RGBColor(0x00, 0xFF, 0x00),
            'blue': RGBColor(0x00, 0x00, 0xFF),
            'yellow': RGBColor(0xFF, 0xFF, 0x00),
            'orange': RGBColor(0xFF, 0xA5, 0x00),
            'purple': RGBColor(0x80, 0x00, 0x80),
            'black': RGBColor(0x00, 0x00, 0x00),
            'white': RGBColor(0xFF, 0xFF, 0xFF),
            'grey': RGBColor(0x80, 0x80, 0x80)
        }

        # Verificar si el color es un valor hexadecimal de 6 dígitos
        if re.match(r'^[0-9A-Fa-f]{6}$', color_name):
            r = int(color_name[0:2], 16)
            g = int(color_name[2:4], 16)
            b = int(color_name[4:6], 16)
            return RGBColor(r, g, b)

        # Si es un nombre de color, buscarlo en el mapa
        return color_mapping.get(color_name.lower(), RGBColor(0x00, 0x00, 0x00))

    def _process_formatting(self, paragraph, text):
        # Diccionario de formatos
        formats = {
            'BOLD': False,
            'ITALIC': False,
            'COLOR': None
        }

        pattern = r'\[(.*?)\]'
        matches = list(re.finditer(pattern, text))
        last_end = 0

        # Procesar los matches encontrados
        for match in matches:
            start, end = match.span()
            tag = match.group(1)

            # Agregar texto antes del tag encontrado
            if last_end < start:
                paragraph.add_run(text[last_end:start])

            # Procesar el tag
            if tag.startswith('COLOR:'):
                color_value = tag.split(':')[1]
                formats['COLOR'] = self._get_color(color_value)
            elif tag == 'BOLD':
                formats['BOLD'] = True
            elif tag == 'ITALIC':
                formats['ITALIC'] = True
            elif tag == 'H1':
                paragraph.style = 'H1'
            elif tag == 'H2':
                paragraph.style = 'H2'

            last_end = end

        # Agregar texto después del último tag
        if last_end < len(text):
            run = paragraph.add_run(text[last_end:])
            if formats['COLOR']:

                run.font.color.rgb = RGBColor.from_string('0000FF')     
            
            if formats['BOLD']:
                run.bold = True
            if formats['ITALIC']:
                run.italic = True

    def _create_table(self, table_data):
        if not table_data:
            return

        rows = [row.split(";") for row in table_data]
        table = self.doc.add_table(rows=len(rows), cols=len(rows[0]))
        table.style = 'Table Grid'

        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                table.cell(i, j).text = cell.strip()

    def _clean_document(self):
        for paragraph in self.doc.paragraphs:
            paragraph.text = self._remove_bracketed_text(paragraph.text)

    def _remove_bracketed_text(self, content):
        return re.sub(r'\[.*?\]', '', content).strip()



from pptx import Presentation
from pptx.dml.color import RGBColor
import re
from pptx.util import Inches

class PptxProcessor:
    def __init__(self):
        self.prs = Presentation()

    def create_presentation(self, content, output_path='output.pptx'):
        self._process_content(content)
        self.prs.save(output_path)
        return output_path

    def _process_content(self, content):
        slides = content.split("[PAGEBREAK]")
        current_table = None
        for slide_content in slides:
            # Si encontramos [TABLE], procesamos una tabla
            if "[TABLE]" in slide_content:
                current_table = []
                slide_content = slide_content.replace("[TABLE]", "").strip()  # Eliminar la marca [TABLE]
            elif "[/TABLE]" in slide_content:
                if current_table:
                    self._create_table(current_table)  # Crear la tabla
                current_table = None
                slide_content = slide_content.replace("[/TABLE]", "").strip()  # Eliminar la marca [/TABLE]
            
            # Procesar el contenido de la diapositiva (con o sin tabla)
            self._create_slide(slide_content.strip(), current_table)

    def _create_table(self, table_data):
        """ Crear la tabla en la diapositiva actual """
        if not table_data:
            return

        # Dividir cada fila de la taula
        rows = [row.split(";") for row in table_data]  # Separar les cel·les per ';'

        # Crear la taula a la diapositiva actual
        slide = self.prs.slides[-1]  # Agafem la diapositiva actual, que és la darrera afegida
        table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(1), Inches(1.5), Inches(8), Inches(5)).table

        # Establir el disseny de la taula (es pot personalitzar més)
        table.style = 'Table Grid'

        # Omplir les cel·les amb les dades
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                table.cell(i, j).text = cell.strip()

    def _create_slide(self, content, current_table=None):
        # Depurar el contenido recibido para entender el problema
        print(f"Contenido recibido para la diapositiva: '{content}'")

        # Comprobar si el contenido está vacío
        if not content.strip():  # Si el contenido está vacío o solo tiene espacios
            print("Advertencia: El contenido está vacío. Se creará una diapositiva vacía.")
            # Si el contenido está vacío, podemos decidir no hacer nada o crear una diapositiva vacía
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Crear diapositiva vacía
            return  # Salir de la función sin procesar más

        # Crear una diapositiva en blanco (layout 5)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Layout en blanco

        # Dividir el contenido en líneas
        lines = content.splitlines()

        # Limpiar las etiquetas [H1] y [/H1] del primer renglón (título)
        title_text = self._remove_bracketed_text(lines[0])  # El primer renglón limpio de etiquetas

        # Colocar el primer renglón como título en el cuadro de título predefinido
        title = slide.shapes.title  # Acceder al cuadro de título
        if title and title_text:
            # Asignar el primer renglón como el título
            title.text = title_text  # El primer renglón limpio se pone como el título

        # Procesar el resto del contenido (si lo hay)
        if len(lines) > 1:
            # Resto del contenido se coloca en un cuadro de texto normal
            self._create_text(slide, "\n".join(lines[1:]), current_table)

    def _remove_bracketed_text(self, content):
        """Eliminar las etiquetas [H1] y [/H1], así como cualquier otra etiqueta similar."""
        # Usamos expresiones regulares para eliminar cualquier etiqueta con formato [TAG] o [/TAG]
        return re.sub(r'\[.*?\]', '', content).strip()  # Eliminar cualquier etiqueta y limpiar los espacios

    def _create_text(self, slide, content, current_table=None):
        # Crear un cuadro de texto para el contenido restante debajo del título
        text_frame = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5)).text_frame
        text_frame.word_wrap = True  # Habilitar ajuste de texto

        # Procesar el contenido del texto
        for line in content.splitlines():
            self._process_line(text_frame, line)

    def _process_line(self, text_frame, text):
        """ Procesar cada línea de texto y aplicar formato según las etiquetas """
        p = text_frame.add_paragraph()
        formats = {
            'BOLD': False,
            'ITALIC': False,
            'COLOR': None
        }

        pattern = r'\[(.*?)\]'  # Buscar las etiquetas como [BOLD], [COLOR], etc.
        matches = list(re.finditer(pattern, text))
        last_end = 0

        for match in matches:
            start, end = match.span()
            tag = match.group(1)

            if last_end < start:
                run = p.add_run()
                run.text = text[last_end:start]

                if formats['BOLD']:
                    run.font.bold = True
                if formats['ITALIC']:
                    run.font.italic = True
                if formats['COLOR']:
                    run.font.color.rgb = formats['COLOR']

            if tag.startswith('COLOR:'):
                color_value = tag.split(':')[1]
                formats['COLOR'] = self._get_color(color_value)
            elif tag == 'BOLD':
                formats['BOLD'] = True
            elif tag == 'ITALIC':
                formats['ITALIC'] = True

            last_end = end

        if last_end < len(text):
            run = p.add_run()
            run.text = text[last_end:]
            if formats['BOLD']:
                run.font.bold = True
            if formats['ITALIC']:
                run.font.italic = True
            if formats['COLOR']:
                run.font.color.rgb = formats['COLOR']

    def _get_color(self, color_name):
        """ Convertir el nombre o el código hexadecimal de un color a RGB """
        color_mapping = {
            'red': RGBColor(255, 0, 0),
            'green': RGBColor(0, 255, 0),
            'blue': RGBColor(0, 0, 255),
            'yellow': RGBColor(255, 255, 0),
            'orange': RGBColor(255, 165, 0),
            'purple': RGBColor(128, 0, 128),
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255),
            'grey': RGBColor(128, 128, 128)
        }

        # Verificar si el color es un valor hexadecimal
        if re.match(r'^[0-9A-Fa-f]{6}$', color_name):
            r = int(color_name[0:2], 16)
            g = int(color_name[2:4], 16)
            b = int(color_name[4:6], 16)
            return RGBColor(r, g, b)

        # Si el color es un nombre, lo buscamos en el mapa
        return color_mapping.get(color_name.lower(), RGBColor(0, 0, 0))

    def _create_table_on_slide(self, slide, table_data):
        """ Crear una tabla en la diapositiva """
        rows = [row.split(";") for row in table_data]  # Usamos ';' para separar las celdas
        table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(1), Inches(1.5), Inches(8), Inches(5)).table

        # Configurar el estilo de la tabla (puedes personalizarlo más)
        table.style = 'Table Grid'

        # Rellenar las celdas con los datos
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                table.cell(i, j).text = cell.strip()

    def _clean_document(self):
        for paragraph in self.prs.paragraphs:
            paragraph.text = self._remove_bracketed_text(paragraph.text)






def generate_index(prompt, file_type, api_key_file='API_KEY.txt', model='mixtral-8x7b-32768'):
    try:
        with open(api_key_file, 'r') as file:
            api_key = file.read().strip()

        client = Groq(api_key=api_key)

        outline_prompt = f"""
        Generate a structured outline for a {file_type} document based on the following prompt:
        - Provide an initial sentence with the summary of the format and lenght of text.
        - Provide a list of sections or pages, each with a title, subtitle (if applicable), and a brief summary of what will be included.
        - Structure must be as this example, splitting pages with //:
        GENERAL PROMPT: Create text for a pptx document with no more than 20 words for each slide.
        PAGE1: TITLE: Exemple de títol CONTENT: Briefing in one sentence
        //
        PAGE2: TITLE: Exemple de títol2 CONTENT: Briefing in one sentence
        //

        Your response must be only the index, no explanation or comments.
        
        Always introduce CONTENT: for every page, do not skip the word CONTENT:
        Here is the prompt for the document outline:
        {prompt}
        """
        print(outline_prompt)

        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": outline_prompt}],
            model=model
        )
        content = chat_completion.choices[0].message.content.strip()
        return content

    except Exception as e:
        raise RuntimeError(f"Error generating index: {e}")


def parse_index(index):
    # Dividir el índice por los delimitadores //
    pages = index.split("//")
    
    # Lista para almacenar las instrucciones procesadas de cada página
    parsed_pages = []

    for page in pages:
        page = page.strip()  # Eliminar posibles espacios extras al principio o final

        if page:  # Ignorar partes vacías
            # Buscar el título y el contenido en cada sección de página
            page_info = {}

            # Buscamos el título
            title_start = page.find("TITLE: ")
            content_start = page.find("CONTENT: ")

            if title_start != -1 and content_start != -1:
                title = page[title_start + len("TITLE: "):content_start].strip()
                content = page[content_start + len("CONTENT: "):].strip()

                # Almacenamos los resultados en el diccionario
                page_info["TITLE"] = title
                page_info["CONTENT"] = content

                # Agregar la página procesada a la lista
                parsed_pages.append(page_info)

    return parsed_pages


def generate_content(index_prompt, index, file_type, api_key_file='API_KEY.txt', model='llama-3.1-70b-versatile'):
    try:
        with open(api_key_file, 'r') as file:
            api_key = file.read().strip()

        client = Groq(api_key=api_key)

        formatting_codes = [
            ("[H1]", "[/H1]", "Title for the main section. It should be large and bold. Do not add it for long sentences."),
            ("[H2]", "[/H2]", "Subtitle for subsections. It should be slightly smaller than [H1] and bold."),
            ("[BOLD]", "[/BOLD]", "Text that should appear in bold."),
            ("[ITALIC]", "[/ITALIC]", "Text that should appear in italics."),
            ("[COLOR:color_name]", "[/COLOR]", "Text that should appear in the specified color. Available colors are red, green, blue, yellow, orange, purple, black, white, grey."),
            ("[TABLE]", "[/TABLE]", "A table. Data inside the table will be organized in rows and columns."),
            ("[PAGEBREAK]", "", "Inserts a page break, starting a new page in the document."),
        ]

        formatting_codes_str = "\n".join([f"- {code[0]}: {code[1]} (Description: {code[2]})" for code in formatting_codes])

        full_prompt = f"""
        Generate structured content for a {file_type} document based on the following prompt:
        - Use the formatting codes provided below ONLY WHERE APPROPIATE. Always add start and end codes.
        Here are the formatting codes and their meanings:
        {formatting_codes_str}
        - For page or slide breaks, simply use [PAGEBREAK].
        - Do not use other formatting codes not provided.
        - For presentations (pptx), include a title for each slide with appropriate formatting (H1).
        - For text documents (docx), include title after a pagebreak and paragraphs with appropriate formatting.
        - You can combine multiple formatting codes, for example: [BOLD][COLOR:red]This is bold and red text[/COLOR][/BOLD]
        - Available colors: red, green, blue, yellow, orange, purple, black, white, grey
        - For tables, use the format:
          [TABLE]
          Header1;Header2
          Row1Col1;Row1Col2
          Row2Col1;Row2Col2
          [/TABLE]
        - Example of formatted output:
          [H1]Exploring the World of Science[/H1]
          [BOLD]Science is a way of learning about the world around us.[/BOLD]
          [ITALIC]This is a bulleted list item.[/ITALIC]
          [COLOR:blue]- Item 1[/COLOR]
          [COLOR:green]- Item 2[/COLOR]
          [PAGEBREAK]
          [TABLE]
          Animals;Habitat
          Lion;Savannah
          Tiger;Forest
          Whale;Ocean
          [/TABLE]
        Do not add any comment rather than the content asked, also do not repeat the prompt or any instruction.
        Always write content in the same language as given instructions.
        Ensure text lenght is appropiate for the type of document and instructions given. 

        The structure or index of the full document is {index}
        From this index you only have to create content for this page (consider full index as contextual):
        {index_prompt}
        """

        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": full_prompt}],
            model=model
        )
        content = chat_completion.choices[0].message.content.strip()
        return content

    except Exception as e:
        raise RuntimeError(f"Error generating content: {e}")



def generate_and_create_file(prompt, file_type, output, model='mixtral-8x7b-32768'):
    try:
        # Primero, generar el índice con los títulos y resúmenes de las páginas
        index = generate_index(prompt, file_type, model=model)  # Suponemos que 'generate_index' genera el índice

        print(index)

        parsed_index = parse_index(index)  # Parseamos el índice en un formato estructurado

        # Crear una lista para almacenar el contenido de cada página
        all_content = []

        # Para cada página en el índice, generamos el contenido
        for page_info in parsed_index:
            page_prompt = page_info["CONTENT"]  # Usamos el contenido de cada página como el prompt
            page_content = generate_content(page_prompt, index, file_type, model=model)  # Generamos el contenido para esta página
            
            # Verificamos el contenido generado
            if not page_content:
                print(f"Advertencia: El contenido de la página '{page_info['TITLE']}' está vacío.")
            else:
                print(f"Contenido de la página '{page_info['TITLE']}': {page_content[:100]}...")  # Mostrar los primeros 100 caracteres para verificar

            all_content.append(page_content)  # Almacenamos el contenido generado

        # Unir el contenido de todas las páginas separadas por [PAGEBREAK]
        combined_content = "[PAGEBREAK]".join(all_content[1:])  # Unir todas excepto la primera
        combined_content = all_content[0] + combined_content  # Añadir la primera página sin [PAGEBREAK]
        
        # Verificamos el contenido combinado antes de pasarlo al procesador
        if not combined_content:
            print("Advertencia: El contenido combinado está vacío.")
        
        # Dependiendo del tipo de archivo, crear el documento adecuado
        if file_type.lower() == 'docx':
            processor = DocxProcessor()
            return processor.create_document(combined_content, output)
        elif file_type.lower() == 'pptx':
            processor = PptxProcessor()
            return processor.create_presentation(combined_content, output)
        else:
            raise ValueError("Tipo de archivo no soportado. Usa 'docx' o 'pptx'.")
    
    except Exception as e:
        raise RuntimeError(f"Error generating and creating file: {e}")