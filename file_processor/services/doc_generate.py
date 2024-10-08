from docx import Document
from docx.shared import RGBColor
from pptx import Presentation
from pptx.util import Inches
from groq import Groq

def apply_color(paragraph, text, color):
    run = paragraph.add_run(text)
    rgb = tuple(map(int, color.split(',')))  # Convertir cadena de color a tupla RGB
    run.font.color.rgb = RGBColor(*rgb)

def add_bold_run(paragraph, text):
    run = paragraph.add_run(text)
    run.bold = True

def add_italic_run(paragraph, text):
    run = paragraph.add_run(text)
    run.italic = True

def create_docx_with_formatting(content, output_path='output.docx'):
    """
    Crea un archivo .docx a partir del contenido generado con formato.

    Args:
        content (str): El contenido generado con códigos de formato.
        output_path (str): La ruta donde se guardará el archivo .docx.
    
    Returns:
        str: La ruta del archivo generado.
    """
    doc = Document()

    # Definir los códigos de formato y sus acciones
    formatting_codes = [
        ("[H1]", "[/H1]", lambda p, text: add_bold_run(p, text)),  # Título principal en negrita
        ("[H2]", "[/H2]", lambda p, text: add_bold_run(p, text)),  # Subtítulo en negrita
        ("[BOLD]", "[/BOLD]", lambda p, text: add_bold_run(p, text)),  # Texto en negrita
        ("[ITALIC]", "[/ITALIC]", lambda p, text: add_italic_run(p, text)),  # Texto en cursiva
        ("[UL]", "[/UL]", lambda p, text: doc.add_paragraph(text, style='ListBullet')),  # Lista con viñetas
        ("[COLOR:", "[/COLOR]", lambda p, text, color: apply_color(p, text, color))  # Texto coloreado
    ]

    for line in content.splitlines():
        p = doc.add_paragraph()  # Crear un nuevo párrafo para cada línea
        while line:
            for start_code, end_code, action in formatting_codes:
                if start_code in line:
                    if start_code.startswith("[COLOR:") and end_code in line:
                        color_start = line.index("[COLOR:") + len("[COLOR:")
                        color_end = line.index("]", color_start)
                        color_name = line[color_start:color_end]
                        colored_text = line.split(f"[COLOR:{color_name}]")[1].split("[/COLOR]")[0]

                        # Aplica el color y elimina la etiqueta
                        action(p, colored_text, color_name)
                        line = line.replace(f"[COLOR:{color_name}]{colored_text}[/COLOR]", "", 1)
                        break

                    elif end_code in line:
                        formatted_text = line.split(start_code)[1].split(end_code)[0]

                        # Aplica la acción y elimina la etiqueta
                        action(p, formatted_text)
                        line = line.replace(f"{start_code}{formatted_text}{end_code}", "", 1)
                        break
            else:
                # Si no hay más códigos, añadimos el texto restante
                p.add_run(line)  # Añadir texto restante
                line = ""  # Finaliza la línea actual

    doc.save(output_path)
    return output_path

def apply_color(paragraph, text, color_name):
    """
    Aplica color al texto dentro de un párrafo.

    Args:
        paragraph (docx.Paragraph): El párrafo donde se añadirá el texto.
        text (str): El texto a colorear.
        color_name (str): El nombre o valor RGB del color.
    """
    try:
        # Supone que color_name viene en formato "R,G,B"
        color = RGBColor(*[int(c) for c in color_name.split(',')])
    except ValueError:
        # Mapeo de colores por nombre si no es un formato RGB
        color_mapping = {
            'red': RGBColor(255, 0, 0),
            'green': RGBColor(0, 255, 0),
            'blue': RGBColor(0, 0, 255),
            'yellow': RGBColor(255, 255, 0),
            'orange': RGBColor(255, 165, 0),
            'purple': RGBColor(128, 0, 128),
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255),
            'grey': RGBColor(128, 128, 128)
        }
        color = color_mapping.get(color_name.lower(), RGBColor(0, 0, 0))

    run = paragraph.add_run(text)
    run.font.color.rgb = color

def create_pptx_with_formatting(content, output_path='output.pptx'):
    """
    Crea un archivo .pptx a partir del contenido generado con formato.

    Args:
        content (str): El contenido generado con códigos de formato.
        output_path (str): La ruta donde se guardará el archivo .pptx.
    
    Returns:
        str: La ruta del archivo generado.
    """
    prs = Presentation()

    # Definir los códigos de formato para PPTX
    formatting_codes = [
        ("[H1]", "[/H1]", lambda s, text: s.shapes.title.text == text),  # Título principal
        ("[H2]", "[/H2]", lambda s, text: s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1)).text == text),  # Subtítulo
        ("[BOLD]", "[/BOLD]", lambda p, text: p.add_run(text).bold == True),  # Texto en negrita
        ("[ITALIC]", "[/ITALIC]", lambda p, text: p.add_run(text).italic == True),  # Texto en cursiva
        ("[UL]", "[/UL]", lambda p, text: p.add_paragraph(text, style='ListBullet')),  # Lista con viñetas
        ("[COLOR:", "[/COLOR]", lambda p, text, color: apply_color_pptx(p, text, color))  # Texto coloreado
        ]

    for line in content.splitlines():
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Slide layout en blanco
        text_frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5)).text_frame
        
        while line:
            for start_code, end_code, action in formatting_codes:
                if start_code in line:
                    # Para etiquetas de color, se necesita un manejo especial para extraer el color
                    if start_code.startswith("[COLOR:") and end_code in line:
                        color_start = line.index("[COLOR:") + len("[COLOR:")
                        color_end = line.index("]", color_start)
                        color_name = line[color_start:color_end]
                        colored_text = line.split(f"[COLOR:{color_name}]")[1].split("[/COLOR]")[0]

                        # Aplica la acción de color y elimina la etiqueta
                        action(text_frame, colored_text, color_name)
                        line = line.replace(f"[COLOR:{color_name}]{colored_text}[/COLOR]", "", 1)
                        break

                    # Para otros formatos estándar (H1, H2, BOLD, ITALIC, UL)
                    elif start_code in line and end_code in line:
                        formatted_text = line.split(start_code)[1].split(end_code)[0]
                        
                        # Aplica la acción y elimina la etiqueta
                        action(text_frame, formatted_text)
                        line = line.replace(f"{start_code}{formatted_text}{end_code}", "", 1)
                        break
            else:
                # Si no hay más códigos, añadimos el texto restante y salimos del bucle
                p = text_frame.add_paragraph()
                p.text = line
                line = ""
    
    prs.save(output_path)
    return output_path

def apply_color_pptx(text_frame, text, color_name):
    """
    Aplica color al texto dentro de un cuadro de texto en PPTX.

    Args:
        text_frame (pptx.text.TextFrame): El cuadro de texto donde se añadirá el texto.
        text (str): El texto a colorear.
        color_name (str): El nombre o valor RGB del color.
    """
    try:
        # Supone que color_name viene en formato "R,G,B"
        color = RGBColor(*[int(c) for c in color_name.split(',')])
    except ValueError:
        # Mapeo de colores por nombre si no es un formato RGB
        color_mapping = {
            'red': RGBColor(255, 0, 0),
            'green': RGBColor(0, 255, 0),
            'blue': RGBColor(0, 0, 255),
            'yellow': RGBColor(255, 255, 0),
            'orange': RGBColor(255, 165, 0),
            'purple': RGBColor(128, 0, 128),
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255),
            'grey': RGBColor(128, 128, 128)
        }
        color = color_mapping.get(color_name.lower(), RGBColor(0, 0, 0))

    p = text_frame.add_paragraph()
    run = p.add_run(text)
    run.font.color.rgb = color


def generate_content(prompt, file_type, api_key_file='API_KEY.txt', model='llama-3.1-70b-versatile'):
    """
    Generates document or presentation content from a prompt using a language model.

    Args:
        prompt (str): The base text for generation.
        api_key_file (str): File containing the API key for the model.
        model (str): The language model to be used for generation.

    Returns:
        str: Generated content in text format with formatting codes.
    """
    try:
        with open(api_key_file, 'r') as file:
            api_key = file.read().strip()

        client = Groq(api_key=api_key)

        # List of formatting codes and their meanings
        formatting_codes = [
            ("[H1]", "[/H1]"),  # Título principal en negrita
            ("[H2]", "[/H2]"),  # Subtítulo en negrita
            ("[BOLD]", "[/BOLD]"),  # Texto en negrita
            ("[ITALIC]", "[/ITALIC]"),  # Texto en cursiva
            ("[UL]", "[/UL]"),  # Lista con viñetas
            ("[OL]", "[/OL]"),  # Lista numerada
            ("[COLOR:", "[/COLOR]")  # Texto coloreado
        ]

        # Formato de la cadena de formato para el prompt
        formatting_codes_str = "\n".join([f"- {code[0]}: {code[1].__doc__}" for code in formatting_codes])

        # Ajusta el prompt para incluir la lista de códigos
        full_prompt = f"""
        Generate a structured content for a {file_type} document based on the following prompt:
        - Use the formatting codes provided below where appropriate. Always add start and end codes.
        Here are the formatting codes and their meanings:
        {formatting_codes_str}
        - Do not use other formatting codes not provided.
        - For presentations (pptx), include a title and bullet points for each slide with appropriate formatting.
        - For text documents (docx), include a title and paragraphs with appropriate formatting.
        - Example of formatted output:
          [H1]Exploring the World of Science[/H1]
          [BOLD]Science is a way of learning about the world around us. [/BOLD]
          [UL]
          [ITALIC]This is a bulleted list item.[/ITALIC]
          [/UL]

        Always write content in same language from given instructions.
        These are the given instructions for the content:
        {prompt}
        """

        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": full_prompt
                }
            ],
            model=model
        )
        content = chat_completion.choices[0].message.content.strip()
        
        # Return the generated content with color mapping
        print(content)
        return content

    except Exception as e:
        raise RuntimeError(f"Error generating content: {e}")
    

def generate_and_create_file(prompt,file_type,output):
    """
    Genera contenido y crea un archivo según el tipo especificado.

    Args:
        file_type (str): Tipo de archivo a generar ('docx' o 'pptx').
        prompt (str): El texto que se usará como prompt para la generación de contenido.
    
    Returns:
        str: La ruta del archivo generado.
    """
    content = generate_content(prompt, file_type)

    if file_type.lower() == 'docx':
        return create_docx_with_formatting(content,output)
    elif file_type.lower() == 'pptx':
        return create_pptx_with_formatting(content,output)
    else:
        raise ValueError("Tipo de archivo no soportado. Usa 'docx' o 'pptx'.")















