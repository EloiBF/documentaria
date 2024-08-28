from pptx import Presentation
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import RGBColor
from pdf2docx import Converter
from transformers import MarianMTModel, MarianTokenizer
import re
import os
from groq import Groq 


# Funciones para asignar un código a cada parte del texto con formato distinto
def generate_numeric_code(counter):
    return f"{counter:06d}"

# Función para no enviar a traducir textos de un color concreto, y mantenerlos como están
def color_to_rgb(color_to_exclude):
    # Verificar si color_hex es None o está vacío
    if color_to_exclude is None or len(color_to_exclude) != 7 or not color_to_exclude.startswith('#'):
        return None
    else:
        # Si pasa la validación, convertir el color hexadecimal a RGB
        color_to_exclude = color_to_exclude.lstrip('#')
        return RGBColor(int(color_to_exclude[:2], 16), int(color_to_exclude[2:4], 16), int(color_to_exclude[4:], 16))


# Funciones para seleccionar textos relevantes y separar el texto en bloques para enviar a traducir
def filtrar_textos_relevantes(textos):  # Quita textos que son solo espacios, numeros, etc
    return {code: texto for code, texto in textos.items() if texto.strip() and not re.fullmatch(r'[\s\W\d]+', texto)}

def unir_textos_fragmentados(textos):  # Hay casos en que una palabra tiene letras con diferente formato, esta formula junta los dos textos para unificar la palabra y le da el primer código
    codigos = sorted(textos.keys())
    nuevo_textos = {}
    buffer_texto = ""
    buffer_codigo = ""
    texto_separador = {}

    for i in range(len(codigos)):
        codigo = codigos[i]
        texto = textos[codigo]

        if buffer_texto:
            if (len(buffer_texto) <= 3 or len(texto) <= 3) and buffer_texto[-1].isalpha() and texto and texto[0].isalpha() and texto[0].islower():
                texto_separador[buffer_codigo] = len(buffer_texto)  # Guardar la longitud del primer texto
                buffer_texto += texto
                textos[codigo] = ""  # Vaciar el texto en el diccionario original
            else:
                nuevo_textos[buffer_codigo] = buffer_texto
                buffer_texto = texto
                buffer_codigo = codigo
        else:
            buffer_texto = texto
            buffer_codigo = codigo

    if buffer_texto:
        nuevo_textos[buffer_codigo] = buffer_texto

    return nuevo_textos, texto_separador

def separar_texto_bloques(textos, max_block_size=200, placeholder='<CD_TR>'):
    def split_text_into_blocks(joined_text, max_block_size):
        blocks = []
        current_block = ""
        current_size = 0

        # Dividir el texto en partes usando el placeholder
        parts = joined_text.split(placeholder)

        for part in parts:
            part_size = len(part) + len(placeholder)  # Considerar el placeholder en el tamaño
            if current_size + part_size > max_block_size:
                blocks.append(current_block.strip(placeholder))
                current_block = part + placeholder
                current_size = len(part) + len(placeholder)
            else:
                current_block += part + placeholder
                current_size += part_size

        if current_block:
            blocks.append(current_block.strip(placeholder))

        return blocks

    # Unir los textos usando el placeholder
    joined_texts = placeholder.join(textos.values())

    if not joined_texts.strip():
        print("Sin texto, manteniendo original")
        return textos

    if re.fullmatch(r'[\s\W\d]+', joined_texts):
        print("No traducir (espacio, símbolo o números), manteniendo original")
        return textos

    bloques = split_text_into_blocks(joined_texts, max_block_size)
    return bloques





# APLICACIÓN DEL MODELO IA DE TRADUCCIÓN -- Entran bloques y salen bloques traducidos    # PODRÍAMOS INCLUIR FUNCIONES DE EMBEDDING AQUÍ
def modelo_traduccion_bloques(bloques, origin_language, destination_language, placeholder='<CD_TR>', numintentos=3):
    
    def verificar_placeholders(placeholder, original_text, translated_text):
        placeholders_originales = re.findall(re.escape(placeholder), original_text)
        placeholders_traducidos = re.findall(re.escape(placeholder), translated_text)
        return len(placeholders_originales) == len(placeholders_traducidos)
    
    def limpiar_placeholders(traduccion):
        return re.sub(r'\s*' + re.escape(placeholder) + r'\s*', placeholder, traduccion)
    
    # Inicializar el cliente de Groq
    with open('API_KEY.txt', 'r') as fichero:
        api_key = fichero.read().strip()

    client = Groq(api_key=api_key)

    bloques_traducidos = []

    for bloque in bloques:
        reintentos = 0
        while reintentos < numintentos:
            try:
                # Crear el mensaje de prompt
                prompt = f"""You are a useful translator specialized in modern Catalan language.
                Translate text ignoring placeholders <CD_TR>, but it is crucial that you maintain them in the same position, even if it is at the start or end of text. 
                Do not add any comments or annotations other than translation. I don't want your feedback. Return only the translation.
                Always use correct catalan gramatical constructions, specially focus in the correct use of ' in articles or pronouns.
                Words that start by capital letter sould also be translated, except if you detect it could be a name.    
                Translate text from {origin_language} to {destination_language}:\n\n{bloque}"""
                
                # Llamar al modelo a través de la API de Groq
                chat_completion = client.chat.completions.create(
                    messages=[
                        {
                            "role": "user",
                            "content": prompt,
                        }
                    ],
                    #model = 'llama3-8b-8192',
                    model = 'llama3-70b-8192' # Aquest model funciona molt bé!
                    #model = 'mixtral-8x7b-32768'
                )
                
                # Obtener la traducción desde la respuesta
                traduccion = chat_completion.choices[0].message.content.strip()
                
                # Limpiar y verificar los placeholders en la traducción
                traduccion = limpiar_placeholders(traduccion)
                
                if verificar_placeholders(placeholder, bloque, traduccion):
                    bloques_traducidos.append(traduccion)
                    print(f'Bloque original: {bloque}')
                    print(f'Bloque traducido: {traduccion}')
                    break
                else:
                    print(f'Bloque original ERROR: {bloque}')
                    print(f'Bloque traducido ERROR: {traduccion}')
                    print('Error en la traducción del bloque, los placeholders no están intactos. Reintentando...')
            except Exception as e:
                print(f'Error en la traducción del bloque: {e}. Reintentando...')
            reintentos += 1
        else:
            print(f'No se pudo traducir correctamente el bloque después de {numintentos} intentos')
            bloques_traducidos.append(bloque)  # Mantener bloque original si falla la traducción

    return bloques_traducidos






def join_blocks(bloques_traducidos, textos_originales, placeholder='<CD_TR>'):
# Inicializar el diccionario traducido
    textos_traducidos = {}
    # Unir los bloques traducidos en un solo texto
    traduccion_completa = placeholder.join(bloques_traducidos)
    # Dividir la traducción completa usando el mismo placeholder
    partes_traducidas = traduccion_completa.split(placeholder)
    # Asegurarse de que el número de partes traducidas coincide con el número de textos originales
    if len(partes_traducidas) != len(textos_originales):
        raise ValueError("El número de bloques traducidos no coincide con el número de textos originales.")
    # Crear un iterador sobre las claves originales
    claves_originales = list(textos_originales.keys())
    # Asignar las partes traducidas al diccionario resultante
    for clave, texto_traducido in zip(claves_originales, partes_traducidas):
        textos_traducidos[clave] = texto_traducido.strip()
    return textos_traducidos


# Limpiamos espacios y puntos que ha añadido de más la traducción
def ajuste_post_traduccion(textos_originales, textos_traducidos):
    # Diccionario para guardar las traducciones ajustadas
    textos_ajustados = {}
    # Iterar sobre cada clave en el diccionario de textos originales
    for key, texto in textos_originales.items():
        traduccion = textos_traducidos.get(key)
        if texto is None or traduccion is None:
            # Si el texto original o la traducción es None, no agregar al diccionario ajustado
            continue
        if texto.strip() == "":
            # Si el texto original está vacío, no agregar al diccionario ajustado
            continue
        leading_spaces = len(texto) - len(texto.lstrip())
        trailing_spaces = len(texto) - len(texto.rstrip())
        # Ajuste de puntos al final de la traducción
        if traduccion.endswith('.') and not texto.endswith('.'):
            traduccion = traduccion.rstrip('.')
        elif traduccion.endswith('..') and not texto.endswith('..'):
            traduccion = traduccion.rstrip('.')
        # Restaurar los espacios al principio y al final
        traduccion = traduccion.rstrip() + ' ' * trailing_spaces
        traduccion = ' ' * leading_spaces + traduccion.lstrip()
        # Guardar la traducción ajustada en el diccionario final
        textos_ajustados[key] = traduccion
    # Filtrar entradas con valores None o vacíos del diccionario ajustado
    textos_ajustados = {k: v for k, v in textos_ajustados.items() if v is not None and v.strip() != ""}
    return textos_ajustados

def separar_palabras_fragmentadas(textos_traducidos, texto_separador, textos_originales): # Una vez tenemos la el diccionario traducido, separamos las palabras unificadas con sus dos códigos, en base al num de caracteres del primer texto.
    textos_traducidos_ok = {}
    nuevos_codigos = {}
    counter = max([int(c) for c in textos_originales.keys()]) + 1

    for codigo, texto_traducido in textos_traducidos.items():
        if codigo in texto_separador:
            longitud = texto_separador[codigo]
            if len(texto_traducido) > longitud:
                primer_texto = texto_traducido[:longitud]
                segundo_texto = texto_traducido[longitud:]

                textos_traducidos_ok[codigo] = primer_texto

                # Encontrar el siguiente código disponible en el texto original
                siguiente_codigo = None
                for cod in sorted(textos_originales.keys()):
                    if cod > codigo and textos_originales[cod].strip() == "":
                        siguiente_codigo = cod
                        break

                # Si no hay un siguiente código disponible, crear uno nuevo
                if siguiente_codigo is None:
                    siguiente_codigo = generate_numeric_code(counter)
                    counter += 1

                textos_traducidos_ok[siguiente_codigo] = segundo_texto
                nuevos_codigos[codigo] = siguiente_codigo
            else:
                textos_traducidos_ok[codigo] = texto_traducido
        else:
            textos_traducidos_ok[codigo] = texto_traducido

    return textos_traducidos_ok



# Lectura/Reemplazo de PPTX
def procesar_ppt(input_path, textos_originales, color_to_exclude, textos_traducidos_final, action): 
    doc = Presentation(input_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for slide in doc.slides:
        for shape in slide.shapes: 

            if shape.has_text_frame: # Acceder a figuras y cuadros de texto
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]

            elif shape.has_table: # Acceder a tablas
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text: # Verificamos que existe la figura
                                    if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                        if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                            code = generate_numeric_code(counter)
                                            counter += 1
                                            if action == "leer":
                                                textos_originales[code] = run.text
                                            elif action == "reemplazar" and code in textos_traducidos_final:
                                                run.text = textos_traducidos_final[code]

            elif shape.has_chart:  # Acceder a gráficos
                chart = shape.chart
                # Título del gráfico
                element = chart.chart_title
                if element.has_text_frame:
                    if element.text_frame.text.strip():
                        code = generate_numeric_code(counter)
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = element.text_frame.text
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            element.text_frame.text = textos_traducidos_final[code]

            elif shape.has_chart:  # Acceder a gráficos
                chart = shape.chart
                for serie in chart:
                    element = serie.point.data_labels
                    if element.has_text_frame:
                        if element.text_frame.text.strip():
                            code = generate_numeric_code(counter)
                            counter += 1
                            if action == "leer":
                                textos_originales[code] = element.text_frame.text
                            elif action == "reemplazar" and code in textos_traducidos_final:
                                element.text_frame.text = textos_traducidos_final[code]

    if action == "leer":
        return textos_originales
    elif action == "reemplazar":
        return doc

# Lectura/Reemplazo DOCX
def procesar_docx(input_path, textos_originales, color_to_exclude, textos_traducidos_final, action): 
    doc = Document(input_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text: # Verificamos que existe la figura
                if run.text.strip(): # Verificamos que tiene texto (no vacía)
                    if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                        code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = run.text
                            print(run.text)
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            run.text = textos_traducidos_final[code]
                            
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter)
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]
    for shape in doc.inline_shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:  # Verificamos que existe la figura
                        if run.text.strip(): # Verificamos que tiene texto (no vacía)
                            if  exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb: # Verificamos el color del texto
                                code = generate_numeric_code(counter)
                                counter += 1
                                if action == "leer":
                                    textos_originales[code] = run.text
                                elif action == "reemplazar" and code in textos_traducidos_final:
                                    run.text = textos_traducidos_final[code]
    if action == "leer":
        return textos_originales
    elif action == "reemplazar":
        return doc

# Lectura/Reemplazo PDF
def procesar_pdf(input_path, textos_originales, color_to_exclude, textos_traducidos_final, action): 
    # Conversor de PDF a WORD
    def pdf_to_word(input_path, word_path):
        # Convertir el PDF a Word
        cv = Converter(input_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()
    # Convertir el PDF a Word
    word_path = input_path.replace('.pdf', '.docx')
    pdf_to_word(input_path, word_path)

    doc = Document(word_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text: # Verificamos que existe la figura
                if run.text.strip(): # Verificamos que tiene texto (no vacía)
                    if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                        code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = run.text
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            run.text = textos_traducidos_final[code]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter)
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]
    for shape in doc.inline_shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:  # Verificamos que existe la figura
                        if run.text.strip(): # Verificamos que tiene texto (no vacía)
                            if  exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb: # Verificamos el color del texto
                                code = generate_numeric_code(counter)
                                counter += 1
                                if action == "leer":
                                    textos_originales[code] = run.text
                                elif action == "reemplazar" and code in textos_traducidos_final:
                                    run.text = textos_traducidos_final[code]
    if action == "leer":
        return textos_originales
    elif action == "reemplazar":
        return doc


# Procesar documento según si es PPT, DOCX o PDF
def procesar_documento(extension, input_path, textos_originales, color_to_exclude, textos_traducidos_final, action):
    if extension == ".pptx":
        return procesar_ppt(input_path, textos_originales, color_to_exclude, textos_traducidos_final, action)
    elif extension == ".docx":
        return procesar_docx(input_path, textos_originales, color_to_exclude, textos_traducidos_final, action)
    elif extension == ".pdf":
        return procesar_pdf(input_path, textos_originales, color_to_exclude, textos_traducidos_final, action) 


# Función final que lee el documento y realiza la traducción. Genera el diccionario original, lo ajusta, recibe el traducido, lo ajusta y reemplaza los textos con los valores del traducido final
def traducir_doc(input_path, output_path, origin_language, destination_language, extension, color_to_exclude):
    print(f"Starting translation process for {input_path}")

    textos_originales = {}
    textos_originales = procesar_documento(extension, input_path, textos_originales, color_to_exclude, textos_traducidos_final=None, action ="leer") 

    print(extension)
    print("Diccionario textos_originales")
    print(textos_originales)

    # Unir textos fragmentados y luego filtrar textos irrelevantes
    textos_para_traducir, texto_separador = unir_textos_fragmentados(textos_originales)
    textos_para_traducir = filtrar_textos_relevantes(textos_para_traducir)

    print("Diccionario textos_para_traducir")
    print(textos_para_traducir)

    # Generación de bloques

    bloques = separar_texto_bloques(textos_para_traducir)

    # Traducción de bloques con el modelo
    bloques_traducidos = modelo_traduccion_bloques(bloques, origin_language, destination_language)

    # Traducir los textos recopilados en bloques --> Obtenemos un diccionario con los textos traducidos
    textos_traducidos = join_blocks(bloques_traducidos,textos_para_traducir)

    print("Diccionario textos_traducidos")
    print(textos_traducidos)

    # Limpiar el texto traducido y separar las palabras fragmentadas
    textos_traducidos_final = ajuste_post_traduccion(textos_para_traducir,textos_traducidos)

    # Separar los textos traducidos --> Generamos el diccionario con los textos que se han traducido y su código, separando las palabras igual que en origen (si venían palabras divididas en varios textos)
    textos_traducidos_final = separar_palabras_fragmentadas(textos_traducidos_final, texto_separador, textos_originales)

    print("Diccionario textos_traducidos_final")
    print(textos_traducidos_final)

    # Generar documento traducido

    doc_traducido = procesar_documento(extension, input_path, textos_originales, color_to_exclude, textos_traducidos_final, action ="reemplazar") 

    # Guardar documento traducido (si es pdf se guarda como docx)
    word_save_path = output_path.replace('.pdf', '.docx')
    if extension ==".pdf":
        doc_traducido.save(word_save_path)
    else:
        doc_traducido.save(output_path)
    print(f'Se ha dejado el documento traducido en la ruta especificada: {output_path}')
