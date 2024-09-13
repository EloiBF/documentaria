import os
import docx
import html2text
from PyPDF2 import PdfReader
from pptx import Presentation
from groq import Groq


def read_document(file_path):
    """Lee el contenido del archivo especificado según su tipo y devuelve el texto y la extensión del archivo."""
    _, file_extension = os.path.splitext(file_path)
    
    if file_extension == '.txt':
        return read_txt(file_path), file_extension
    
    elif file_extension == '.docx':
        return read_docx(file_path), file_extension
    
    elif file_extension == '.pdf':
        return read_pdf(file_path), file_extension
    
    elif file_extension == '.pptx':
        return read_pptx(file_path), file_extension
    
    elif file_extension == '.html' or file_extension == '.htm':
        return read_html(file_path), file_extension
    
    else:
        raise ValueError(f"Tipo de archivo no soportado: {file_extension}")

def read_txt(file_path):
    """Lee el contenido de un archivo .txt."""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_docx(file_path):
    """Lee el contenido de un archivo .docx."""
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def read_pdf(file_path):
    """Lee el contenido de un archivo .pdf."""
    reader = PdfReader(file_path)
    full_text = []
    for page in reader.pages:
        full_text.append(page.extract_text())
    return '\n'.join(full_text)

def read_pptx(file_path):
    """Lee el contenido de un archivo .pptx."""
    presentation = Presentation(file_path)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)

def read_html(file_path):
    """Lee el contenido de un archivo .html o .htm."""
    with open(file_path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    return html2text.html2text(html_content)

def split_text(text, max_length=5000):
    """Divide el texto en bloques, asegurándose de que cada bloque no supere max_length caracteres."""
    blocks = []
    start = 0
    while start < len(text):
        end = start + max_length
        if end < len(text):
            # Intentamos cortar en un punto lógico, buscando el último punto o salto de línea antes del límite de caracteres.
            cut_index = text.rfind('.', start, end)
            if cut_index == -1:  # Si no hay un punto, cortamos directamente en el límite.
                cut_index = end
            else:
                cut_index += 1  # Incluir el punto en el bloque.
        else:
            cut_index = len(text)
        
        blocks.append(text[start:cut_index].strip())
        start = cut_index
    
    return blocks




def summarize_text(texto, file_type=None, model='llama-3.1-70b-versatile', api_key_file='API_KEY.txt', num_words=None , summary_language="auto", add_prompt=""):
    try:
        with open(api_key_file, 'r') as fichero:
            api_key = fichero.read().strip()
        client = Groq(api_key=api_key)

        if summary_language == "auto":
            summary_language == "the same language as given text"

        if file_type == None:
            prompt = f"""Explain me the content of the text in a summarized way strictly following this rules:
             - Focus on the most important topics in text
             - Use aprox {num_words} words for the whole summary
             - Do not add any extra comment, annotation or introduction to the response. Print only the summarized text.
             - Summary must be generated in {summary_language}
            Additional rules that must be followed:
            {add_prompt}
             Text to summarize is:{texto}"""

        elif file_type in ['.pptx', '.docx', '.pdf','.txt', '.html']:
            prompt = f"""Explain me the content of the text in a summarized way strictly following this rules:
             - Focus on the most important topics in text
             - Use aprox {num_words} words for the whole summary
             - Do not add any extra comment, annotation or introduction to the response. Print only the summarized text.
             - Ignore codes (_CDTR_00000) as if they were not in the text.
             - Summary must be generated in {summary_language}
            Additional rules that must be followed:
            {add_prompt}
             
             Text to summarize is:{texto}"""

        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": prompt,

                }
            ],
            model=model
        )
        traduccion = chat_completion.choices[0].message.content.strip()
        return traduccion

    except Exception as e:
        raise RuntimeError(f"Error during translation: {e}")





def generate_summary(text, num_words, summary_language, file_type, add_prompt, old_summary=""):
    """Genera un resumen del texto utilizando el modelo especificado y maneja bloques si es necesario."""
    blocks = split_text(text)
    summary = ""
    
    for i, block in enumerate(blocks):
        if i > 0:
            block = old_summary + block  # Añade el contexto del bloque anterior.
        partial_summary = summarize_text(block, file_type=file_type, num_words=num_words, summary_language=summary_language, add_prompt=add_prompt)
        summary += partial_summary + "\n"
    
    return summary.strip()

def resumir_doc(input_path, num_words, summary_language, add_prompt):
    """Función principal para resumir un documento y guardar el resumen en un archivo TXT."""
    try:
        # Leer el documento y obtener el tipo de archivo
        text, file_type = read_document(input_path)
        
        # Generar el resumen 
        summary = generate_summary(text, num_words, summary_language, file_type, add_prompt, old_summary="Resumen anterior: ")
        
        return summary  # Retorna el resumen 
        
    except Exception as e:
        print(f"Ocurrió un error: {e}")
        return None, None  # Retorna None en caso de error
