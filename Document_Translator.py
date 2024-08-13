print('Ejecuta')
import os
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import RGBColor
from pdf2docx import Converter
from transformers import M2M100ForConditionalGeneration, M2M100Tokenizer
import time
import re
from pathlib import Path

# VARIABLES:
nombre_fichero = 'PRUEBA 1.pptx' # Debe incluir la extensión al final (.pptx, .docx o .pdf)
origin_language = "es"
destination_language = "ca"
color_to_exclude = "#FF00FF"
# Configuración Embedding: Idiomas disponibles: CAS, CAT, FRA, ALE, ING
idioma_original = "CAS"
idioma_traducir = "CAT"
delimitadores = r'[.!?\n:]'
texto_minimo = 5

# Rutas de documentos a traducir y devolución
input_path = f'test/in/{nombre_fichero}'
output_path = f'test/out/Traducido_{destination_language}_{nombre_fichero}'
extension = Path(nombre_fichero).suffix

# Cargar el modelo y el tokenizador
model_path = "models/facebook_m2m100_418M"
tokenizer = M2M100Tokenizer.from_pretrained(model_path)
model = M2M100ForConditionalGeneration.from_pretrained(model_path)

# Funciones para asignar un código a cada parte del texto con formato distinto
def generate_numeric_code(counter):
    return f"{counter:06d}"

# Función para no enviar a traducir textos de un color concreto, y mantenerlos como están
def color_to_rgb(color_hex): 
    color_hex = color_hex.lstrip('#')
    return RGBColor(int(color_hex[:2], 16), int(color_hex[2:4], 16), int(color_hex[4:], 16))

# Funciones para seleccionar textos relevantes y separar el texto en bloques para enviar a traducir
def filtrar_textos_relevantes(textos):  # Quita textos que son solo espacios, numeros, etc
    return {code: texto for code, texto in textos.items() if texto.strip() and not re.fullmatch(r'[\s\W\d]+', texto)}

def unir_textos_fragmentados(textos):  # Hay casos en que una palabra tiene letras con diferente formato, esta formula junta los dos textos para unificar la palabra y le da el primer código
    codigos = sorted(textos.keys())
    nuevo_textos = {}
    buffer_texto = ""
    buffer_codigo = ""
    texto_separador = {}

    for i in range(len(codigos)):
        codigo = codigos[i]
        texto = textos[codigo]

        if buffer_texto:
            if (len(buffer_texto) <= 3 or len(texto) <= 3) and buffer_texto[-1].isalpha() and texto and texto[0].isalpha() and texto[0].islower():
                texto_separador[buffer_codigo] = len(buffer_texto)  # Guardar la longitud del primer texto
                buffer_texto += texto
                textos[codigo] = ""  # Vaciar el texto en el diccionario original
            else:
                nuevo_textos[buffer_codigo] = buffer_texto
                buffer_texto = texto
                buffer_codigo = codigo
        else:
            buffer_texto = texto
            buffer_codigo = codigo

    if buffer_texto:
        nuevo_textos[buffer_codigo] = buffer_texto

    return nuevo_textos, texto_separador

def separar_texto_bloques(textos, max_block_size):
    def split_text_into_blocks(textos, max_block_size): # Separamos el texto en bloques, asegurando no cortar ningún code ni texto por la mitad
        blocks = []
        current_block = ""
        current_size = 0

        pattern = re.compile(r'\(_CDTR\d{6}\)')

        parts = re.split(f'({pattern.pattern})', textos)
        for i in range(0, len(parts), 2):
            code = parts[i]
            if i + 1 < len(parts):
                text_part = parts[i + 1]
            else:
                text_part = ""

            if current_size + len(code) + len(text_part) > max_block_size:
                blocks.append(current_block)
                current_block = code + text_part
                current_size = len(code) + len(text_part)
            else:
                current_block += code + text_part
                current_size += len(code) + len(text_part)

        if current_block:
            blocks.append(current_block)

        return blocks
      
    joined_texts = "".join([f"(_CDTR_{code}){text}" for code, text in textos.items()])

    if not joined_texts.strip():
        print("Sin texto, manteniendo original")
        return {code: text for code, text in textos.items()}

    if re.fullmatch(r'[\s\W\d]+', joined_texts):
        print("No traducir (espacio, simbolo o numeros), manteniendo original")
        return {code: text for code, text in textos.items()}

    bloques =split_text_into_blocks(joined_texts, max_block_size) # Generamos los bloques
    return bloques




# PODRÍAMOS INCLUIR FUNCIONES DE EMBEDDING AQUÍ


# APLICACIÓN DEL MODELO IA DE TRADUCCIÓN -- Entran bloques y salen bloques traducidos
def modelo_traduccion_bloques(bloques, origin_language, destination_language):
    
    bloques_traducidos = []

    for bloque in bloques:
        # Configurar el idioma fuente y destino
        tokenizer.src_lang = origin_language
        tokenizer.tgt_lang = destination_language

        # Preparar el contexto (más estructurado)
        contexto = f"Eres un traductor de textos que ignora los códigos CDTR y los mantiene una vez traducido el texto. Texto a traducir de español a catalan: {bloque}"
        
        # Tokenizar el texto con contexto
        tokens = tokenizer(contexto, return_tensors='pt', padding=True, truncation=True)
        
        # Generar la traducción
        translated_tokens = model.generate(**tokens)
        traduccion = tokenizer.batch_decode(translated_tokens, skip_special_tokens=True)[0]
        
        bloques_traducidos.append(traduccion)
        print(f'Bloque original: {bloque}')
        print(f'Bloque traducido: {traduccion}')

    return bloques_traducidos




def join_blocks(bloques_traducidos):
    traduccion_completa = "".join(bloques_traducidos)
    translated_texts = re.findall(r"\(_CDTR_([A-Za-z0-9]{6})\)(.*?)(?=\(_CDTR_[A-Za-z0-9]{6}\)|$)", traduccion_completa, re.DOTALL)
    return {code: text for code, text in translated_texts}

# Funciones para limpiar el texto traducido
def eliminar_codigos(diccionario):
    # Expresión regular para eliminar códigos con o sin paréntesis, con una o dos barras bajas, y con cualquier cantidad de cifras o letras
    patron = r"\(?_?CDTR[_A-Za-z0-9]{1,}\)?|\(?__?CDTR[_A-Za-z0-9]{1,}\)?|CDTR[_A-Za-z0-9]{1,}"
    # Aplica la eliminación de códigos a cada valor en el diccionario
    return {key: re.sub(patron, "", value) for key, value in diccionario.items()}

# Limpiamos espacios y puntos que ha añadido de más la traducción
def ajuste_post_traduccion(textos_originales, textos_traducidos):
    # Diccionario para guardar las traducciones ajustadas
    textos_ajustados = {}
    # Iterar sobre cada clave en el diccionario de textos originales
    for key, texto in textos_originales.items():
        traduccion = textos_traducidos.get(key)
        if texto is None or traduccion is None:
            # Si el texto original o la traducción es None, no agregar al diccionario ajustado
            continue
        if texto.strip() == "":
            # Si el texto original está vacío, no agregar al diccionario ajustado
            continue
        leading_spaces = len(texto) - len(texto.lstrip())
        trailing_spaces = len(texto) - len(texto.rstrip())
        # Ajuste de puntos al final de la traducción
        if traduccion.endswith('.') and not texto.endswith('.'):
            traduccion = traduccion.rstrip('.')
        elif traduccion.endswith('..') and not texto.endswith('..'):
            traduccion = traduccion.rstrip('.')
        # Restaurar los espacios al principio y al final
        traduccion = traduccion.rstrip() + ' ' * trailing_spaces
        traduccion = ' ' * leading_spaces + traduccion.lstrip()
        # Guardar la traducción ajustada en el diccionario final
        textos_ajustados[key] = traduccion
    # Filtrar entradas con valores None o vacíos del diccionario ajustado
    textos_ajustados = {k: v for k, v in textos_ajustados.items() if v is not None and v.strip() != ""}
    return textos_ajustados

def separar_palabras_fragmentadas(textos_traducidos, texto_separador, textos_originales): # Una vez tenemos la el diccionario traducido, separamos las palabras unificadas con sus dos códigos, en base al num de caracteres del primer texto.
    textos_traducidos_ok = {}
    nuevos_codigos = {}
    counter = max([int(c) for c in textos_originales.keys()]) + 1

    for codigo, texto_traducido in textos_traducidos.items():
        if codigo in texto_separador:
            longitud = texto_separador[codigo]
            if len(texto_traducido) > longitud:
                primer_texto = texto_traducido[:longitud]
                segundo_texto = texto_traducido[longitud:]

                textos_traducidos_ok[codigo] = primer_texto

                # Encontrar el siguiente código disponible en el texto original
                siguiente_codigo = None
                for cod in sorted(textos_originales.keys()):
                    if cod > codigo and textos_originales[cod].strip() == "":
                        siguiente_codigo = cod
                        break

                # Si no hay un siguiente código disponible, crear uno nuevo
                if siguiente_codigo is None:
                    siguiente_codigo = generate_numeric_code(counter)
                    counter += 1

                textos_traducidos_ok[siguiente_codigo] = segundo_texto
                nuevos_codigos[codigo] = siguiente_codigo
            else:
                textos_traducidos_ok[codigo] = texto_traducido
        else:
            textos_traducidos_ok[codigo] = texto_traducido

    return textos_traducidos_ok



# Lectura/Reemplazo de PPTX
def procesar_ppt(input_path, textos_originales, textos_traducidos_final, action): 
    doc = Presentation(input_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for slide in doc.slides:
        for shape in slide.shapes: 

            if shape.has_text_frame: # Acceder a figuras y cuadros de texto
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]

            elif shape.has_table: # Acceder a tablas
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text: # Verificamos que existe la figura
                                    if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                        if run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                            code = generate_numeric_code(counter)
                                            counter += 1
                                            if action == "leer":
                                                textos_originales[code] = run.text
                                            elif action == "reemplazar" and code in textos_traducidos_final:
                                                run.text = textos_traducidos_final[code]

            elif shape.has_chart:  # Acceder a gráficos
                chart = shape.chart
                # Título del gráfico
                element = chart.chart_title
                if element.has_text_frame:
                    if element.text_frame.text.strip():
                        code = generate_numeric_code(counter)
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = element.text_frame.text
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            element.text_frame.text = textos_traducidos_final[code]

            elif shape.has_chart:  # Acceder a gráficos
                chart = shape.chart
                for serie in chart:
                    element = serie.point.data_labels
                    if element.has_text_frame:
                        if element.text_frame.text.strip():
                            code = generate_numeric_code(counter)
                            counter += 1
                            if action == "leer":
                                textos_originales[code] = element.text_frame.text
                            elif action == "reemplazar" and code in textos_traducidos_final:
                                element.text_frame.text = textos_traducidos_final[code]

    if action == "leer":
        return textos_originales
    elif action == "reemplazar":
        return doc

# Lectura/Reemplazo DOCX
def procesar_docx(input_path, textos_originales, textos_traducidos_final, action): 
    doc = Document(input_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text: # Verificamos que existe la figura
                if run.text.strip(): # Verificamos que tiene texto (no vacía)
                    if run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                        code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = run.text
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            run.text = textos_traducidos_final[code]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter)
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]
    for shape in doc.inline_shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:  # Verificamos que existe la figura
                        if run.text.strip(): # Verificamos que tiene texto (no vacía)
                            if  run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb: # Verificamos el color del texto
                                code = generate_numeric_code(counter)
                                counter += 1
                                if action == "leer":
                                    textos_originales[code] = run.text
                                elif action == "reemplazar" and code in textos_traducidos_final:
                                    run.text = textos_traducidos_final[code]
    if action == "leer":
        return textos_originales
    elif action == "reemplazar":
        return doc

# Lectura/Reemplazo PDF
def procesar_pdf(input_path, textos_originales, textos_traducidos_final, action): 
    # Conversor de PDF a WORD
    def pdf_to_word(input_path, word_path):
        # Convertir el PDF a Word
        cv = Converter(input_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()
    # Convertir el PDF a Word
    word_path = input_path.replace('.pdf', '.docx')
    pdf_to_word(input_path, word_path)

    doc = Document(word_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text: # Verificamos que existe la figura
                if run.text.strip(): # Verificamos que tiene texto (no vacía)
                    if run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                        code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = run.text
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            run.text = textos_traducidos_final[code]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter)
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]
    for shape in doc.inline_shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:  # Verificamos que existe la figura
                        if run.text.strip(): # Verificamos que tiene texto (no vacía)
                            if  run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb: # Verificamos el color del texto
                                code = generate_numeric_code(counter)
                                counter += 1
                                if action == "leer":
                                    textos_originales[code] = run.text
                                elif action == "reemplazar" and code in textos_traducidos_final:
                                    run.text = textos_traducidos_final[code]
    if action == "leer":
        return textos_originales
    elif action == "reemplazar":
        return doc


# Procesar documento según si es PPT, DOCX o PDF
def procesar_documento(input_path, textos_originales, textos_traducidos_final, action):
    if extension == ".pptx":
        return procesar_ppt(input_path, textos_originales, textos_traducidos_final, action)
    elif extension == ".docx":
        return procesar_docx(input_path, textos_originales, textos_traducidos_final, action)
    elif extension == ".pdf":
        return procesar_pdf(input_path, textos_originales, textos_traducidos_final, action) 








# Función final que lee el documento y realiza la traducción. Genera el diccionario original, lo ajusta, recibe el traducido, lo ajusta y reemplaza los textos con los valores del traducido final
def traducir_doc(input_path, output_path):
    
    textos_originales = {}
    textos_originales = procesar_documento(input_path, textos_originales, textos_traducidos_final=None, action ="leer") 

    print("Diccionario textos_originales")
    print(textos_originales)

    # Unir textos fragmentados y luego filtrar textos irrelevantes
    textos_para_traducir, texto_separador = unir_textos_fragmentados(textos_originales)
    textos_para_traducir = filtrar_textos_relevantes(textos_para_traducir)

    print("Diccionario textos_para_traducir")
    print(textos_para_traducir)

    # Generación de bloques

    bloques = separar_texto_bloques(textos_para_traducir,max_block_size=500)

    # Traducción de bloques con el modelo
    bloques_traducidos = modelo_traduccion_bloques(bloques, origin_language, destination_language)

    # Traducir los textos recopilados en bloques --> Obtenemos un diccionario con los textos traducidos
    textos_traducidos = join_blocks(bloques_traducidos)

    print("Diccionario textos_traducidos")
    print(textos_traducidos)

    # Limpiar el texto traducido y separar las palabras fragmentadas
    textos_traducidos_final = eliminar_codigos(textos_traducidos)
    textos_traducidos_final = ajuste_post_traduccion(textos_para_traducir,textos_traducidos_final)

    # Separar los textos traducidos --> Generamos el diccionario con los textos que se han traducido y su código, separando las palabras igual que en origen (si venían palabras divididas en varios textos)
    textos_traducidos_final = separar_palabras_fragmentadas(textos_traducidos_final, texto_separador, textos_originales)

    print("Diccionario textos_traducidos_final")
    print(textos_traducidos_final)

    # Generar documento traducido

    doc_traducido = procesar_documento(input_path, textos_originales, textos_traducidos_final, action ="reemplazar") 

    # Guardar documento traducido (si es pdf se guarda como docx)
    word_save_path = output_path.replace('.pdf', '.docx')
    if extension ==".pdf":
        doc_traducido.save(word_save_path)
    else:
        doc_traducido.save(output_path)
    print(f'Se ha dejado el documento traducido en la ruta especificada: {output_path}')

# Ejecutar la traducción del documento
traducir_doc(input_path, output_path)

