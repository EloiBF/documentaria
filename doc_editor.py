from pptx import Presentation
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import RGBColor
from pdf2docx import Converter
import re
from bs4 import BeautifulSoup
import chardet
from groq import Groq
import zipfile
from lxml import etree
from docx import Document
from io import BytesIO
import shutil


# Funciones para asignar un código a cada parte del texto con formato distinto
def generate_numeric_code(counter):
    return f"{counter:06d}"

# Función para no enviar a traducir textos de un color concreto, y mantenerlos como están
def color_to_rgb(color_to_exclude):
    # Verificar si color_hex es None o está vacío
    if color_to_exclude is None or len(color_to_exclude) != 7 or not color_to_exclude.startswith('#'):
        return None
    else:
        # Si pasa la validación, convertir el color hexadecimal a RGB
        color_to_exclude = color_to_exclude.lstrip('#')
        return RGBColor(int(color_to_exclude[:2], 16), int(color_to_exclude[2:4], 16), int(color_to_exclude[4:], 16))


# Funciones para seleccionar textos relevantes y separar el texto en bloques para enviar a traducir
def filtrar_textos_relevantes(textos):  # Quita textos que son solo espacios, numeros, etc
    return {code: texto for code, texto in textos.items() if texto.strip() and not re.fullmatch(r'[\s\W\d]+', texto)}

def unir_textos_fragmentados(textos):  # Hay casos en que una palabra tiene letras con diferente formato, esta formula junta los dos textos para unificar la palabra y le da el primer código
    codigos = sorted(textos.keys())
    nuevo_textos = {}
    buffer_texto = ""
    buffer_codigo = ""
    texto_separador = {}

    for i in range(len(codigos)):
        codigo = codigos[i]
        texto = textos[codigo]

        if buffer_texto:
            if (len(buffer_texto) <= 3 or len(texto) <= 3) and buffer_texto[-1].isalpha() and texto and texto[0].isalpha() and texto[0].islower():
                texto_separador[buffer_codigo] = len(buffer_texto)  # Guardar la longitud del primer texto
                buffer_texto += texto
                textos[codigo] = ""  # Vaciar el texto en el diccionario original
            else:
                nuevo_textos[buffer_codigo] = buffer_texto
                buffer_texto = texto
                buffer_codigo = codigo
        else:
            buffer_texto = texto
            buffer_codigo = codigo

    if buffer_texto:
        nuevo_textos[buffer_codigo] = buffer_texto

    return nuevo_textos, texto_separador

def separar_texto_bloques(textos, max_keys_per_block=10):
    """Separar el texto en bloques basados en un número máximo de claves por bloque."""
    
    def split_text_into_blocks(textos, max_keys_per_block):
        blocks = []
        current_block = []
        current_keys = 0

        # Crear una lista de pares (código, texto) a partir del diccionario de textos
        items = list(textos.items())
        
        for code, text in items:
            if current_keys >= max_keys_per_block:
                # Si el número máximo de claves por bloque se ha alcanzado, añadir el bloque actual y empezar uno nuevo
                blocks.append(current_block)
                current_block = []
                current_keys = 0
            
            # Añadir el par (código, texto) al bloque actual
            current_block.append(f"(_CDTR_{code}){text}")
            current_keys += 1

        # Añadir el último bloque si no está vacío
        if current_block:
            blocks.append(current_block)

        return blocks

    # Convertir el diccionario de textos en bloques basados en el número máximo de claves por bloque
    bloques = split_text_into_blocks(textos, max_keys_per_block)
    
    # Convertir cada lista de bloques en una cadena de texto
    bloques_texto = ["".join(bloque) for bloque in bloques]
    
    return bloques_texto

# APLICACIÓN DEL MODELO IA DE TRADUCCIÓN -- Entran bloques y salen bloques traducidos    # PODRÍAMOS INCLUIR FUNCIONES DE EMBEDDING AQUÍ
# Se usan las funciones de model_translator.py, que se puede usar para traducir cualquier texto

def verificar_codigos(original, traducido):
    """Verifica que todos los códigos en el texto original estén presentes en el texto traducido."""
    patron_codigo = r'\(_CDTR_\d{6}\)'
    codigos_originales = re.findall(patron_codigo, original)
    codigos_traducidos = re.findall(patron_codigo, traducido)
    codigos_originales_set = set(codigos_originales)
    codigos_traducidos_set = set(codigos_traducidos)
    codigos_faltantes = codigos_originales_set - codigos_traducidos_set

    if codigos_faltantes:
        print(f'[ERROR] Códigos faltantes en la traducción: {codigos_faltantes}')
        return False
        
    return True


def modelo_edición_bloques(bloques, add_prompt, extension, numintentos=10):
    bloques_traducidos = []

    for bloque in bloques:
        reintentos = 0
        while reintentos < numintentos:
            try:
                traduccion = prompt_text(bloque, add_prompt, extension)
                
                if verificar_codigos(bloque, traduccion):
                    print(f'Bloque original: {bloque}')
                    print(f'Bloque traducido: {traduccion}')
                    bloques_traducidos.append(traduccion)
                    break
                else:
                    print(f'[ERROR] Bloque original: {bloque}')
                    print(f'[ERROR] Bloque traducido: {traduccion}')
                    raise ValueError("Error por traducción no válida")
            
            except Exception as e:
                print(f'Error en la traducción del bloque. {e}. Reintentando...')
                reintentos += 1
        else:
            print(f'No se pudo traducir correctamente el bloque después de {numintentos} intentos')
            bloques_traducidos.append(bloque)
    
    return bloques_traducidos


def join_blocks(bloques_traducidos):
    traduccion_completa = "".join(bloques_traducidos)
    translated_texts = re.findall(r"\(_CDTR_([A-Za-z0-9]{6})\)(.*?)(?=\(_CDTR_[A-Za-z0-9]{6}\)|$)", traduccion_completa, re.DOTALL)
    return {code: text for code, text in translated_texts}


# Funciones para limpiar el texto traducido
def eliminar_codigos(diccionario):
    # Expresión regular para eliminar códigos con o sin paréntesis, con una o dos barras bajas, y con cualquier cantidad de cifras o letras
    patron = r"\(?_?CDTR[_A-Za-z0-9]{1,}\)?|\(?__?CDTR[_A-Za-z0-9]{1,}\)?|CDTR[_A-Za-z0-9]{1,}"
    # Aplica la eliminación de códigos a cada valor en el diccionario
    return {key: re.sub(patron, "", value) for key, value in diccionario.items()}


# Limpiamos espacios y puntos que ha añadido de más la traducción
def ajuste_post_traduccion(textos_originales, textos_traducidos):
    # Diccionario para guardar las traducciones ajustadas
    textos_ajustados = {}
    # Iterar sobre cada clave en el diccionario de textos originales
    for key, texto in textos_originales.items():
        traduccion = textos_traducidos.get(key)
        if texto is None or traduccion is None:
            # Si el texto original o la traducción es None, no agregar al diccionario ajustado
            continue
        if texto.strip() == "":
            # Si el texto original está vacío, no agregar al diccionario ajustado
            continue
        leading_spaces = len(texto) - len(texto.lstrip())
        trailing_spaces = len(texto) - len(texto.rstrip())
        # Ajuste de puntos al final de la traducción
        if traduccion.endswith('.') and not texto.endswith('.'):
            traduccion = traduccion.rstrip('.')
        elif traduccion.endswith('..') and not texto.endswith('..'):
            traduccion = traduccion.rstrip('.')
        # Restaurar los espacios al principio y al final
        traduccion = traduccion.rstrip() + ' ' * trailing_spaces
        traduccion = ' ' * leading_spaces + traduccion.lstrip()
        # Guardar la traducción ajustada en el diccionario final
        textos_ajustados[key] = traduccion
    # Filtrar entradas con valores None o vacíos del diccionario ajustado
    textos_ajustados = {k: v for k, v in textos_ajustados.items() if v is not None and v.strip() != ""}
    return textos_ajustados

def separar_palabras_fragmentadas(textos_traducidos, texto_separador, textos_originales): # Una vez tenemos la el diccionario traducido, separamos las palabras unificadas con sus dos códigos, en base al num de caracteres del primer texto.
    textos_traducidos_ok = {}
    nuevos_codigos = {}
    counter = max([int(c) for c in textos_originales.keys()]) + 1

    for codigo, texto_traducido in textos_traducidos.items():
        if codigo in texto_separador:
            longitud = texto_separador[codigo]
            if len(texto_traducido) > longitud:
                primer_texto = texto_traducido[:longitud]
                segundo_texto = texto_traducido[longitud:]

                textos_traducidos_ok[codigo] = primer_texto

                # Encontrar el siguiente código disponible en el texto original
                siguiente_codigo = None
                for cod in sorted(textos_originales.keys()):
                    if cod > codigo and textos_originales[cod].strip() == "":
                        siguiente_codigo = cod
                        break

                # Si no hay un siguiente código disponible, crear uno nuevo
                if siguiente_codigo is None:
                    siguiente_codigo = generate_numeric_code(counter)
                    counter += 1

                textos_traducidos_ok[siguiente_codigo] = segundo_texto
                nuevos_codigos[codigo] = siguiente_codigo
            else:
                textos_traducidos_ok[codigo] = texto_traducido
        else:
            textos_traducidos_ok[codigo] = texto_traducido

    return textos_traducidos_ok



# Lectura/Reemplazo de PPTX
def procesar_ppt(input_path,output_path, textos_originales, color_to_exclude, textos_traducidos_final, action): 
    doc = Presentation(input_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for slide in doc.slides:
        for shape in slide.shapes: 

            if shape.has_text_frame: # Acceder a figuras y cuadros de texto
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]

            elif shape.has_table: # Acceder a tablas
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text: # Verificamos que existe la figura
                                    if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                        if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                            code = generate_numeric_code(counter)
                                            counter += 1
                                            if action == "leer":
                                                textos_originales[code] = run.text
                                            elif action == "reemplazar" and code in textos_traducidos_final:
                                                run.text = textos_traducidos_final[code]

            elif shape.has_chart:  # Acceder a gráficos
                chart = shape.chart
                # Título del gráfico
                element = chart.chart_title
                if element.has_text_frame:
                    if element.text_frame.text.strip():
                        code = generate_numeric_code(counter)
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = element.text_frame.text
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            element.text_frame.text = textos_traducidos_final[code]

            elif shape.has_chart:  # Acceder a gráficos
                chart = shape.chart
                for serie in chart:
                    element = serie.point.data_labels
                    if element.has_text_frame:
                        if element.text_frame.text.strip():
                            code = generate_numeric_code(counter)
                            counter += 1
                            if action == "leer":
                                textos_originales[code] = element.text_frame.text
                            elif action == "reemplazar" and code in textos_traducidos_final:
                                element.text_frame.text = textos_traducidos_final[code]

    if action == "leer":
        return textos_originales
    elif action == "reemplazar":
        return doc.save(output_path)

# Lectura/Reemplazo DOCX
def procesar_docx(input_path,output_path, textos_originales, color_to_exclude, textos_traducidos_final, action): 
    doc = Document(input_path)

    exclude_color_rgb = color_to_rgb(color_to_exclude)

    textos_originales = {}
    counter = 1  # Inicializar el contador desde 1

    # Recopilar textos y asignar códigos - Genera el diccionario "textos_originales"
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text: # Verificamos que existe la figura
                if run.text.strip(): # Verificamos que tiene texto (no vacía)
                    if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                        code = generate_numeric_code(counter) # Volvemos a generar un código para cada texto para cruzar el formato con la traducción
                        counter += 1
                        if action == "leer":
                            textos_originales[code] = run.text
                            print(run.text)
                        elif action == "reemplazar" and code in textos_traducidos_final:
                            run.text = textos_traducidos_final[code]
                            
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text: # Verificamos que existe la figura
                            if run.text.strip(): # Verificamos que tiene texto (no vacía)
                                if exclude_color_rgb == None or run.font.color is None or not hasattr(run.font.color, 'rgb') or run.font.color.rgb != exclude_color_rgb:
                                    code = generate_numeric_code(counter)
                                    counter += 1
                                    if action == "leer":
                                        textos_originales[code] = run.text
                                    elif action == "reemplazar" and code in textos_traducidos_final:
                                        run.text = textos_traducidos_final[code]
    # Procesar shapes (cuadros de texto y formas) en el XML
    with zipfile.ZipFile(input_path, 'r') as docx_zip:
        # Extraemos el archivo `document.xml` para modificarlo
        with docx_zip.open('word/document.xml') as document_xml:
            tree = etree.parse(document_xml)
            namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

            # Buscar cuadros de texto (w:txbxContent) y otras figuras (formas)
            cuadros_texto = tree.xpath('//w:txbxContent//w:t', namespaces=namespaces)
            for element in cuadros_texto:
                if action == "leer":
                    code = generate_numeric_code(counter)
                    counter += 1
                    textos_originales[code] = element.text
                elif action == "reemplazar":
                    code = generate_numeric_code(counter)
                    counter += 1
                    if code in textos_traducidos_final:
                        element.text = textos_traducidos_final[code]

            # Buscar otras formas incrustadas (shapes)
            formas = tree.xpath('//w:drawTextBox//w:t', namespaces=namespaces)
            for element in formas:
                if action == "leer":
                    code = generate_numeric_code(counter)
                    counter += 1
                    textos_originales[code] = element.text
                elif action == "reemplazar":
                    code = generate_numeric_code(counter)
                    counter += 1
                    if code in textos_traducidos_final:
                        element.text = textos_traducidos_final[code]

    # Si estamos en modo reemplazar, creamos un nuevo archivo .docx con los cambios
    if action == "reemplazar":
        # Crear un archivo temporal
        temp_zip_path = output_path + '_temp.zip'
        with zipfile.ZipFile(input_path, 'r') as docx_zip:
            with zipfile.ZipFile(temp_zip_path, 'w') as temp_zip:
                # Copiamos todos los archivos originales, excepto `document.xml`
                for item in docx_zip.infolist():
                    if item.filename != 'word/document.xml':
                        temp_zip.writestr(item, docx_zip.read(item.filename))
                # Escribimos el nuevo `document.xml` modificado
                with BytesIO() as buffer:
                    tree.write(buffer, xml_declaration=True, encoding='UTF-8')
                    buffer.seek(0)
                    temp_zip.writestr('word/document.xml', buffer.read())

        # Renombrar el archivo temporal como el archivo de salida final
        shutil.move(temp_zip_path, output_path)

    return textos_originales if action == "leer" else None

# Lectura/Reemplazo PDF
def procesar_pdf(input_path, output_path, textos_originales, color_to_exclude, textos_traducidos_final, action): 
    # Conversor de PDF a WORD
    def pdf_to_word(input_path, word_path):
        # Convertir el PDF a Word
        cv = Converter(input_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()
    # Convertir el PDF a Word
    word_path = input_path.replace('.pdf', '.docx')
    pdf_to_word(input_path, word_path)

    return procesar_docx(word_path,output_path, textos_originales, color_to_exclude, textos_traducidos_final, action)

def procesar_txt(input_path, output_path, action, textos_traducidos_final=None):
    """Leer un archivo TXT, detectar su codificación, y dividir el texto en bloques de 10 palabras."""
    textos_originales = {}
    
    # Detectar la codificación del archivo
    with open(input_path, 'rb') as file:
        raw_data = file.read()
        result = chardet.detect(raw_data)
        encoding = result['encoding']
    
    # Leer el archivo usando la codificación detectada
    with open(input_path, 'r', encoding=encoding) as file:
        text = file.read()
    
    words = text.split()
    num_words = len(words)
    counter = 1
    
    # Dividir el texto en bloques de 10 palabras
    for i in range(0, num_words, 10):
        block = ' '.join(words[i:i+10])
        code = generate_numeric_code(counter)
        counter += 1
        if action == "leer":
            textos_originales[code] = block
        elif action == "reemplazar" and textos_traducidos_final and code in textos_traducidos_final:
            text = text.replace(block, textos_traducidos_final[code])
    
    # Guardar el archivo modificado si la acción es reemplazar
    if action == "reemplazar":
        with open(output_path, 'w', encoding=encoding) as file:
            file.write(text)
    
    # Retornar el diccionario de textos originales si la acción es leer
    return textos_originales if action == "leer" else None


def procesar_html(input_path, output_path, action, textos_traducidos_final=None):
    """Leer un archivo HTML, extraer el texto, y reemplazarlo si es necesario."""
    textos_originales = {}

    # Detectar la codificación del archivo
    with open(input_path, 'rb') as file:
        raw_data = file.read()
        result = chardet.detect(raw_data)
        encoding = result['encoding']

    # Leer el archivo HTML usando la codificación detectada
    with open(input_path, 'r', encoding=encoding) as file:
        soup = BeautifulSoup(file, 'html.parser')

    counter = 1

    # Iterar sobre los elementos del HTML y procesar el texto
    for element in soup.find_all(text=True):
        text = element.strip()
        if text:
            code = generate_numeric_code(counter)
            counter += 1
            if action == "leer":
                textos_originales[code] = text
            elif action == "reemplazar":
                # Buscar el texto traducido usando el código generado
                if code in textos_traducidos_final:
                    new_text = textos_traducidos_final[code]
                    element.replace_with(new_text)

    # Guardar el archivo modificado si la acción es reemplazar
    if action == "reemplazar":
        with open(output_path, 'w', encoding=encoding) as file:
            file.write(str(soup))

    # Retornar el diccionario de textos originales si la acción es leer
    return textos_originales if action == "leer" else None


# Procesar documento según si es PPT, DOCX o PDF
def procesar_documento(extension, input_path ,output_path, textos_originales, color_to_exclude, textos_traducidos_final, action):
    if extension == ".pptx":
        return procesar_ppt(input_path, output_path, textos_originales, color_to_exclude, textos_traducidos_final, action)
    elif extension == ".docx":
        return procesar_docx(input_path, output_path, textos_originales, color_to_exclude, textos_traducidos_final, action)
    elif extension == ".pdf":
        return procesar_pdf(input_path, output_path, textos_originales, color_to_exclude, textos_traducidos_final, action) 
    elif extension in ['.txt']:
        return procesar_txt(input_path, output_path, action, textos_traducidos_final) 
    elif extension in ['.html']:
        return procesar_html(input_path, output_path, action, textos_traducidos_final) 


# Funció genèrica per traduir amb la IA, li passes un text i retorna la traducció

def prompt_text(texto, add_prompt, file_type, model='llama3-70b-8192', api_key_file='API_KEY.txt'):
    try:
        # Inicializa el cliente de Groq
        with open(api_key_file, 'r') as fichero:
            api_key = fichero.read().strip()
        client = Groq(api_key=api_key)

                # Genera el prompt base para la traducción según el tipo de archivo
        if file_type == None:
            prompt = f"""
        You are a multilingual text editor. Follow this rules:
        - Provide only the edited text without any additional comments or annotations. I don't want your feedback, only the pure edition, do not introduce"
        - Keep similar text lenght when editing. 
        - Do not duplicate or translate text if it is not asked, even if instructions are in other language.
        Additional rules that must be followed:
        {add_prompt}
        Text to edit:
        {texto}
        """

        if file_type in ['.pptx', '.docx', '.pdf','.txt', '.html']:
            prompt = f"""
        You are a multilingual text editor. Follow this rules:
        - Your task is to edit text while strictly preserving codes (_CDTR_00000) in their exact positions, including at the start or end of the text. 
        - Provide only the asked text without any additional comments or annotations. I don't want your feedback, only the pure edition, do not introduce"
        - Keep similar text lenght when editing.
        - Do not duplicate or translate text if it is not asked, even if instructions are in other language.

        Additional rules that must be followed:
        {add_prompt}
        Text to edit:
        {texto}
        """
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        # Llama a la API de Groq para editar el texto
        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": prompt,
                }
            ],
            model=model
        )

        traduccion = chat_completion.choices[0].message.content.strip()
        return traduccion
    
    except Exception as e:
        raise RuntimeError(f"Error during translation: {e}")


# Función que aplica todo

def editar_doc(input_path, output_path, extension, color_to_exclude, add_prompt):
    print(f"Starting translation process for {input_path}")

    # Procesar el documento para obtener textos originales
    textos_originales = procesar_documento(extension, input_path, output_path, {}, color_to_exclude, textos_traducidos_final=None, action="leer") 

    print("Diccionario textos_originales")
    print(textos_originales)

    # Unir textos fragmentados y luego filtrar textos irrelevantes
    textos_para_traducir, texto_separador = unir_textos_fragmentados(textos_originales)
    textos_para_traducir = filtrar_textos_relevantes(textos_para_traducir)

    print("Diccionario textos_para_traducir")
    print(textos_para_traducir)

    # Inicializar variable para almacenar textos traducidos
    textos_traducidos = {}

    # Generación de bloques
    bloques = separar_texto_bloques(textos_para_traducir)

    # Traducción de bloques con el modelo
    bloques_traducidos = modelo_edición_bloques(bloques, add_prompt, extension)

    # Traducir los textos recopilados en bloques
    textos_traducidos = join_blocks(bloques_traducidos)

    print("Diccionario textos_traducidos")
    print(textos_traducidos)

    # Limpiar el texto traducido y separar las palabras fragmentadas
    textos_traducidos_final = eliminar_codigos(textos_traducidos)
    textos_traducidos_final = ajuste_post_traduccion(textos_traducidos_final, textos_traducidos)

    # Separar los textos traducidos --> Generamos el diccionario con los textos traducidos y su código
    textos_traducidos_final = separar_palabras_fragmentadas(textos_traducidos_final, texto_separador, textos_originales)

    print("Diccionario textos_traducidos_final")
    print(textos_traducidos_final)

    # Generar documento traducido
    procesar_documento(extension, input_path, output_path, textos_originales, color_to_exclude, textos_traducidos_final, action="reemplazar") 
    
    print(f'Se ha dejado el documento traducido en la ruta especificada: {output_path}')