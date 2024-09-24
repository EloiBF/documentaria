import os
import docx
import html2text
from PyPDF2 import PdfReader
from pptx import Presentation
from groq import Groq
from lxml import etree
import zipfile
from io import BytesIO
import chardet

# Script per llegir tot el text d'un document.
# Serveix per quan no és necessari tornar a generar el document, es llegeix tot de cop i se li passa a la IA blocs de molts caràcters. Útil per fer un resum o extreure informació concreta.


def read_document(file_path):
    """Lee el contenido del archivo especificado según su tipo y devuelve el texto y la extensión del archivo."""
    _, file_extension = os.path.splitext(file_path)
    
    if file_extension == '.txt':
        return read_txt(file_path), file_extension
    
    elif file_extension == '.docx':
        return read_docx(file_path), file_extension
    
    elif file_extension == '.pdf':
        return read_pdf(file_path), file_extension
    
    elif file_extension == '.pptx':
        return read_pptx(file_path), file_extension
    
    elif file_extension == '.html' or file_extension == '.htm':
        return read_html(file_path), file_extension
    
    else:
        raise ValueError(f"Tipo de archivo no soportado: {file_extension}")

def read_txt(file_path):
    """Lee el contenido de un archivo .txt, detectando automáticamente la codificación."""
    encodings = ['utf-8', 'latin-1', 'windows-1252']
    content = None
    
    # Detectar la codificación usando chardet
    try:
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            result = chardet.detect(raw_data)
            encoding_detected = result['encoding']
            print(f"Codificación detectada automáticamente: {encoding_detected}")

            # Intentamos leer con la codificación detectada
            try:
                content = raw_data.decode(encoding_detected)
                print(f"Archivo leído con éxito usando la codificación detectada: {encoding_detected}")
            except (UnicodeDecodeError, TypeError):
                print(f"No se pudo decodificar con la codificación detectada: {encoding_detected}. Intentando otras codificaciones...")

    except Exception as e:
        print(f"Error al leer el archivo en modo binario: {e}")

    # Si la detección automática falla o no es precisa, probamos otras codificaciones
    if content is None:
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read()
                    print(f"Archivo leído con éxito usando la codificación: {encoding}")
                    break  # Si tiene éxito, salimos del loop
            except UnicodeDecodeError:
                print(f"Error al leer el archivo con la codificación: {encoding}")
            except Exception as e:
                print(f"Ocurrió otro error al leer el archivo: {e}")

    return content

def read_docx(file_path):
    """Lee todo el contenido de un archivo .docx, incluyendo cuadros de texto y formas."""
    # Leer el contenido de los párrafos y tablas usando python-docx
    doc = docx.Document(file_path)
    full_text = []

    # Leer texto de párrafos
    for para in doc.paragraphs:
        full_text.append(para.text)

    # Leer texto de tablas
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    full_text.append(para.text)

    # Leer el contenido del XML para capturar cuadros de texto y formas
    with zipfile.ZipFile(file_path, 'r') as docx_zip:
        # Leer el archivo `document.xml`
        with docx_zip.open('word/document.xml') as document_xml:
            tree = etree.parse(document_xml)
            namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

            # Extraer el texto de los cuadros de texto (w:txbxContent)
            cuadros_texto = tree.xpath('//w:txbxContent//w:t', namespaces=namespaces)
            for element in cuadros_texto:
                full_text.append(element.text)

            # Extraer el texto de las formas (w:drawTextBox)
            formas = tree.xpath('//w:drawTextBox//w:t', namespaces=namespaces)
            for element in formas:
                full_text.append(element.text)

    return '\n'.join(filter(None, full_text))  # Filtramos elementos vacíos y unimos con saltos de línea

def read_pdf(file_path):
    """Lee el contenido de un archivo .pdf."""
    reader = PdfReader(file_path)
    full_text = []
    for page in reader.pages:
        full_text.append(page.extract_text())
    return '\n'.join(full_text)

def read_pptx(file_path):
    """Lee el contenido de un archivo .pptx."""
    presentation = Presentation(file_path)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)

def read_html(file_path):
    """Lee el contenido de un archivo .html o .htm."""
    with open(file_path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    return html2text.html2text(html_content)

def split_text(text, max_length=5000):
    """Divide el texto en bloques, asegurándose de que cada bloque no supere max_length caracteres."""
    blocks = []
    start = 0
    while start < len(text):
        end = start + max_length
        if end < len(text):
            # Intentamos cortar en un punto lógico, buscando el último punto o salto de línea antes del límite de caracteres.
            cut_index = text.rfind('.', start, end)
            if cut_index == -1:  # Si no hay un punto, cortamos directamente en el límite.
                cut_index = end
            else:
                cut_index += 1  # Incluir el punto en el bloque.
        else:
            cut_index = len(text)
        
        blocks.append(text[start:cut_index].strip())
        start = cut_index
    
    return blocks
